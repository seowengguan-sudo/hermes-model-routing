#!/usr/bin/env python3
"""
Local Multi-Format Document Reader Agent
=========================================
Reads PDF, CSV, Excel, Word, PPT, RTF, images, and text files entirely locally.
No external LLM/API calls required.

Usage:
    python3 doc_reader.py <file_path> [--extract-tables] [--render-pages] [--ocr-images]

Supported formats:
  - PDF: .pdf (text + tables + page rendering)
  - Excel: .xls, .xlsx, .xlsm (data + formulas + tables)
  - CSV: .csv (structured data)
  - Word: .docx, .doc (text + tables)
  - PowerPoint: .pptx, .ppt (text + slide notes)
  - RTF: .rtf (text)
  - Images: .png, .jpg, .jpeg, .gif, .bmp, .tiff, .webp (render + analyze)
  - Text: .txt, .md, .log, .json, .yaml, .yml, .xml, .html, .htm, .py, .js, etc.
"""

from __future__ import annotations

import argparse
import csv
import io
import json
import os
import re
import sys
import traceback
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

# Add project root to path for Hermes tools
PROJECT_ROOT = "/opt/hermes"
sys.path.insert(0, PROJECT_ROOT)
sys.path.insert(0, "/opt/data/.venv-docreader/lib/python3.13/site-packages")


# ── Data Structures ──────────────────────────────────────────────────────────

@dataclass
class DocumentResult:
    """Result of reading a document."""
    file_path: str
    file_type: str
    file_size: int
    page_count: int | None = None
    text: str = ""
    tables: list[list[list[str]]] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)
    rendered_pages: list[str] = field(default_factory=list)  # paths to rendered images
    error: str | None = None


@dataclass
class FileInfo:
    """Information about a file."""
    path: Path
    name: str
    ext: str
    size: int
    mime_type: str


# ── File Type Detection ──────────────────────────────────────────────────────

PDF_EXTS = {'.pdf'}
EXCEL_EXTS = {'.xls', '.xlsx', '.xlsm', '.xlsb'}
CSV_EXTS = {'.csv'}
WORD_EXTS = {'.docx', '.doc'}
PPT_EXTS = {'.pptx', '.ppt'}
RTF_EXTS = {'.rtf'}
IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.webp'}
TEXT_EXTS = {'.txt', '.md', '.log', '.json', '.yaml', '.yml', '.xml', '.html', '.htm',
             '.py', '.js', '.ts', '.tsx', '.jsx', '.css', '.scss', '.sql',
             '.sh', '.bash', '.zsh', '.conf', '.cfg', '.ini', '.toml',
             '.csv', '.tsv'}

# MIME type mapping (basic, for common cases)
MIME_MAP = {
    '.pdf': 'application/pdf',
    '.xls': 'application/vnd.ms-excel',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.xlsm': 'application/vnd.ms-excel.sheets.macroEnabled.12',
    '.csv': 'text/csv',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.doc': 'application/msword',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.ppt': 'application/vnd.ms-powerpoint',
    '.rtf': 'application/rtf',
    '.txt': 'text/plain',
    '.md': 'text/markdown',
}


def get_file_info(path: str) -> FileInfo:
    """Get information about a file."""
    p = Path(path).resolve()
    if not p.exists():
        raise FileNotFoundError(f"File not found: {path}")
    if not p.is_file():
        raise ValueError(f"Not a file: {path}")
    
    ext = p.suffix.lower()
    size = p.stat().st_size
    mime = MIME_MAP.get(ext, 'application/octet-stream')
    
    return FileInfo(
        path=p,
        name=p.name,
        ext=ext,
        size=size,
        mime_type=mime,
    )


def categorize_file(ext: str) -> str:
    """Categorize a file by its extension."""
    ext = ext.lower()
    if ext in PDF_EXTS:
        return 'pdf'
    elif ext in EXCEL_EXTS:
        return 'excel'
    elif ext in CSV_EXTS:
        return 'csv'
    elif ext in WORD_EXTS:
        return 'word'
    elif ext in PPT_EXTS:
        return 'powerpoint'
    elif ext in RTF_EXTS:
        return 'rtf'
    elif ext in IMAGE_EXTS:
        return 'image'
    elif ext in TEXT_EXTS:
        return 'text'
    else:
        return 'unknown'


# ── PDF Reader ───────────────────────────────────────────────────────────────

def read_pdf(file_path: Path, extract_tables: bool = True, render_pages: bool = False) -> DocumentResult:
    """Read a PDF file using pypdf and pymupdf."""
    result = DocumentResult(
        file_path=str(file_path),
        file_type='pdf',
        file_size=file_path.stat().st_size,
    )
    
    # Extract text with pypdf
    try:
        import pypdf
        reader = pypdf.PdfReader(str(file_path))
        result.metadata = {
            'title': reader.metadata.title if reader.metadata else None,
            'author': reader.metadata.author if reader.metadata else None,
            'subject': reader.metadata.subject if reader.metadata else None,
            'creator': reader.metadata.creator if reader.metadata else None,
            'producer': reader.metadata.producer if reader.metadata else None,
            'page_count': len(reader.pages),
        }
        result.page_count = len(reader.pages)
        
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            result.text += f"\n--- Page {i+1} ---\n{text}\n"
    except Exception as e:
        result.error = f"pypdf error: {e}"
    
    # Fallback: try PyMuPDF for text and tables if pypdf failed
    if not result.text.strip():
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(file_path))
            result.page_count = doc.page_count
            for i, page in enumerate(doc):
                text = page.get_text()
                result.text += f"\n--- Page {i+1} ---\n{text}\n"
            
            # Extract tables with PyMuPDF
            if extract_tables:
                for i, page in enumerate(doc):
                    tables = page.find_tables()
                    for table in tables.tables:
                        table_data = table.extract()
                        if table_data:
                            result.tables.append(table_data)
            
            # Render pages if requested
            if render_pages:
                import tempfile
                for i, page in enumerate(doc):
                    mat = fitz.Matrix(2, 2)  # 2x scale for better quality
                    pix = page.get_pixmap(matrix=mat)
                    img_path = str(file_path.parent / f"{file_path.stem}_page_{i+1}.png")
                    pix.save(img_path)
                    result.rendered_pages.append(img_path)
            
            result.error = None  # Clear previous error if PyMuPDF succeeded
        except Exception as e:
            if not result.error:
                result.error = f"pymupdf error: {e}"
    
    # Extract tables with pdfplumber as additional fallback
    if extract_tables and not result.tables:
        try:
            import pdfplumber
            with pdfplumber.open(str(file_path)) as pdf:
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            result.tables.append(table)
        except Exception:
            pass
    
    return result


# ── Excel Reader ─────────────────────────────────────────────────────────────

def read_excel(file_path: Path, extract_tables: bool = True) -> DocumentResult:
    """Read an Excel file using openpyxl (xls/xlsx) or xlrd (xls)."""
    result = DocumentResult(
        file_path=str(file_path),
        file_type='excel',
        file_size=file_path.stat().st_size,
    )
    
    ext = file_path.suffix.lower()
    
    try:
        if ext in {'.xlsx', '.xlsm', '.xlsb'}:
            import openpyxl
            # Use read-only mode for large files
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
        elif ext == '.xls':
            import xlrd
            wb = xlrd.open_workbook(str(file_path), on_demand=True)
        else:
            # Fallback: try openpyxl
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
        
        result.page_count = len(wb.sheetnames) if hasattr(wb, 'sheetnames') else 1
        result.metadata = {'sheet_names': wb.sheetnames if hasattr(wb, 'sheetnames') else []}
        
        if ext == '.xls':
            # xlrd path
            for sheet_name in wb.sheet_names():
                sheet = wb.sheet_by_name(sheet_name)
                result.text += f"\n--- Sheet: {sheet_name} ---\n"
                rows = []
                for row_idx in range(min(sheet.nrows, 1000)):
                    row = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                    rows.append(row)
                    result.text += "\t".join(row) + "\n"
                if extract_tables:
                    result.tables.append(rows)
        else:
            # openpyxl path
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                result.text += f"\n--- Sheet: {sheet_name} ---\n"
                rows = []
                for row in ws.iter_rows(values_only=True, max_row=1000):
                    row_strs = [str(v) if v is not None else "" for v in row]
                    rows.append(row_strs)
                    result.text += "\t".join(row_strs) + "\n"
                if extract_tables:
                    result.tables.append(rows)
        
        if hasattr(wb, 'close'):
            wb.close()
    except Exception as e:
        result.error = f"Excel reader error: {e}"
    
    return result


# ── CSV Reader ───────────────────────────────────────────────────────────────

def read_csv(file_path: Path, extract_tables: bool = True) -> DocumentResult:
    """Read a CSV file."""
    result = DocumentResult(
        file_path=str(file_path),
        file_type='csv',
        file_size=file_path.stat().st_size,
    )
    
    try:
        # Try UTF-8 first, then fallback
        content = None
        for encoding in ['utf-8', 'utf-8-sig', 'cp1252', 'latin-1']:
            try:
                content = file_path.read_text(encoding=encoding)
                break
            except UnicodeDecodeError:
                continue
        
        if content is None:
            content = file_path.read_text(errors='replace')
        
        # Auto-detect delimiter
        first_line = content.split('\n')[0]
        delimiter = ','
        if '\t' in first_line:
            delimiter = '\t'
        elif ';' in first_line and first_line.count(';') > first_line.count(','):
            delimiter = ';'
        elif '|' in first_line:
            delimiter = '|'
        
        reader = csv.reader(io.StringIO(content), delimiter=delimiter)
        rows = []
        for row in reader:
            rows.append(row)
            result.text += delimiter.join(str(v) for v in row) + "\n"
        
        result.page_count = len(rows)
        result.metadata = {
            'row_count': len(rows),
            'column_count': len(rows[0]) if rows else 0,
            'delimiter': delimiter,
        }
        
        if extract_tables and rows:
            result.tables.append(rows)
    except Exception as e:
        result.error = f"CSV reader error: {e}"
    
    return result


# ── Word Reader ──────────────────────────────────────────────────────────────

def read_word(file_path: Path, extract_tables: bool = True) -> DocumentResult:
    """Read a Word document (docx/doc)."""
    result = DocumentResult(
        file_path=str(file_path),
        file_type='word',
        file_size=file_path.stat().st_size,
    )
    
    ext = file_path.suffix.lower()
    
    try:
        if ext == '.docx':
            import docx
            doc = docx.Document(str(file_path))
            result.page_count = len(doc.paragraphs)
            result.metadata = {
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables),
                'style': doc.styles.element.xml[:200] if doc.styles else None,
            }
            
            for para in doc.paragraphs:
                result.text += para.text + "\n"
            
            # Extract tables
            if extract_tables:
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        table_data.append([cell.text for cell in row.cells])
                    result.tables.append(table_data)
        else:
            # .doc - try to read as RTF or use striprtf
            try:
                from striprtf.striprtf import rtf_to_text
                content = file_path.read_bytes()
                text = rtf_to_text(content.decode('utf-8', errors='replace'))
                result.text = text
            except Exception:
                # Fallback: try reading raw text
                result.text = file_path.read_text(errors='replace')
    except Exception as e:
        result.error = f"Word reader error: {e}"
    
    return result


# ── PowerPoint Reader ────────────────────────────────────────────────────────

def read_pptx(file_path: Path, extract_tables: bool = True) -> DocumentResult:
    """Read a PowerPoint presentation."""
    result = DocumentResult(
        file_path=str(file_path),
        file_type='powerpoint',
        file_size=file_path.stat().st_size,
    )
    
    try:
        import pptx
        prs = pptx.Presentation(str(file_path))
        result.page_count = len(prs.slides)
        result.metadata = {
            'slide_count': len(prs.slides),
        }
        
        for i, slide in enumerate(prs.slides):
            result.text += f"\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    result.text += shape.text + "\n"
            
            # Extract notes
            if slide.has_notes_slide:
                notes = slide.notes_slide
                if hasattr(notes, 'notes_text_frame'):
                    result.text += f"\n[Notes: {notes.notes_text_frame.text}]\n"
    except Exception as e:
        result.error = f"PowerPoint reader error: {e}"
    
    return result


# ── RTF Reader ───────────────────────────────────────────────────────────────

def read_rtf(file_path: Path) -> DocumentResult:
    """Read an RTF file."""
    result = DocumentResult(
        file_path=str(file_path),
        file_type='rtf',
        file_size=file_path.stat().st_size,
    )
    
    try:
        from striprtf.striprtf import rtf_to_text
        content = file_path.read_text(encoding='utf-8', errors='replace')
        text = rtf_to_text(content)
        result.text = text
        result.page_count = 1
    except Exception as e:
        result.error = f"RTF reader error: {e}"
    
    return result


# ── Image Reader ─────────────────────────────────────────────────────────────

def read_image(file_path: Path) -> DocumentResult:
    """Read an image file and extract basic info + render."""
    result = DocumentResult(
        file_path=str(file_path),
        file_type='image',
        file_size=file_path.stat().st_size,
    )
    
    try:
        from PIL import Image
        img = Image.open(str(file_path))
        result.metadata = {
            'format': img.format,
            'mode': img.mode,
            'width': img.width,
            'height': img.height,
            'info': {k: str(v)[:200] for k, v in img.info.items()},
        }
        result.page_count = 1
        result.text = f"Image: {result.metadata['format']} {result.metadata['width']}x{result.metadata['height']} ({result.metadata['mode']})"
        
        # If we have the local vision model, we can use vision_analyze
        # For now, just extract the image metadata
        # The agent can use vision_analyze separately for OCR/content analysis
        
    except Exception as e:
        result.error = f"Image reader error: {e}"
    
    return result


# ── Text Reader ──────────────────────────────────────────────────────────────

def read_text(file_path: Path) -> DocumentResult:
    """Read a plain text file."""
    result = DocumentResult(
        file_path=str(file_path),
        file_type='text',
        file_size=file_path.stat().st_size,
    )
    
    try:
        # Try UTF-8 first, then fallback
        content = None
        for encoding in ['utf-8', 'utf-8-sig', 'cp1252', 'latin-1']:
            try:
                content = file_path.read_text(encoding=encoding)
                result.metadata['encoding'] = encoding
                break
            except UnicodeDecodeError:
                continue
        
        if content is None:
            content = file_path.read_text(errors='replace')
            result.metadata['encoding'] = 'utf-8 (with replacement)'
        
        result.text = content
        result.page_count = content.count('\n') + 1
    except Exception as e:
        result.error = f"Text reader error: {e}"
    
    return result


# ── Main Document Reader ─────────────────────────────────────────────────────

def read_document(file_path: str, extract_tables: bool = True, render_pages: bool = False) -> DocumentResult:
    """
    Read any supported document type entirely locally.
    
    Args:
        file_path: Path to the file to read
        extract_tables: Whether to extract tables from the document
        render_pages: Whether to render PDF/image pages to PNG
    
    Returns:
        DocumentResult with text, tables, metadata, and rendered page paths
    """
    info = get_file_info(file_path)
    category = categorize_file(info.ext)
    
    result = DocumentResult(
        file_path=str(info.path),
        file_type=category,
        file_size=info.size,
        metadata={
            'filename': info.name,
            'extension': info.ext,
            'mime_type': info.mime_type,
            'size_bytes': info.size,
        }
    )
    
    if category == 'pdf':
        return read_pdf(info.path, extract_tables=extract_tables, render_pages=render_pages)
    elif category == 'excel':
        return read_excel(info.path, extract_tables=extract_tables)
    elif category == 'csv':
        return read_csv(info.path, extract_tables=extract_tables)
    elif category == 'word':
        return read_word(info.path, extract_tables=extract_tables)
    elif category == 'powerpoint':
        return read_pptx(info.path, extract_tables=extract_tables)
    elif category == 'rtf':
        return read_rtf(info.path)
    elif category == 'image':
        return read_image(info.path)
    elif category == 'text':
        return read_text(info.path)
    else:
        result.error = f"Unsupported file type: {info.ext}"
        return result


# ── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Local Multi-Format Document Reader Agent",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument('file', help='Path to the file to read')
    parser.add_argument('--extract-tables', action='store_true', default=True,
                       help='Extract tables from the document (default: True)')
    parser.add_argument('--no-tables', dest='extract_tables', action='store_false',
                       help='Skip table extraction')
    parser.add_argument('--render-pages', action='store_true',
                       help='Render PDF/image pages to PNG files')
    parser.add_argument('--json', action='store_true',
                       help='Output as JSON')
    parser.add_argument('--text-only', action='store_true',
                       help='Output only extracted text')
    
    args = parser.parse_args()
    
    try:
        result = read_document(args.file, extract_tables=args.extract_tables, render_pages=args.render_pages)
    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)
    
    if args.json:
        output = {
            'file_path': result.file_path,
            'file_type': result.file_type,
            'file_size': result.file_size,
            'page_count': result.page_count,
            'text': result.text,
            'tables': result.tables,
            'metadata': result.metadata,
            'rendered_pages': result.rendered_pages,
            'error': result.error,
        }
        print(json.dumps(output, indent=2, ensure_ascii=False))
    elif args.text_only:
        print(result.text)
    else:
        # Formatted output
        print(f"📄 {result.file_path}")
        print(f"   Type: {result.file_type} | Size: {result.file_size:,} bytes")
        if result.page_count:
            print(f"   Pages: {result.page_count}")
        if result.metadata:
            for k, v in result.metadata.items():
                if k in ('filename', 'extension', 'mime_type', 'size_bytes', 'encoding'):
                    continue
                print(f"   {k}: {v}")
        if result.tables:
            print(f"   Tables found: {len(result.tables)}")
        if result.rendered_pages:
            print(f"   Rendered pages: {len(result.rendered_pages)} files")
        if result.error:
            print(f"   ⚠ Error: {result.error}")
        print(f"\n{'─'*60}\n")
        print(result.text[:5000] if len(result.text) > 5000 else result.text)
        if len(result.text) > 5000:
            print(f"\n... ({len(result.text) - 5000} more characters)")
    
    return 0 if not result.error else 1


if __name__ == "__main__":
    sys.exit(main())
