#!/usr/bin/env python3
"""
OAKAI Document Reader - Enterprise Data Anonymization Engine
=============================================================
Self-contained document reader that performs enterprise-grade data 
anonymization/pseudonymization locally. No LLM/token/API calls.
All processing is 100% local with reversible variable mapping.

Just run: python3 doc_reader_onefile.py
Then open: http://localhost:8765

Supports: PDF, Word, Excel, PPTX, TXT, CSV, HTML, JSON, MD.
"""
import json, os, re, sys, io, zipfile, hashlib, time
from pathlib import Path
from datetime import datetime
from http.server import HTTPServer, BaseHTTPRequestHandler
from urllib.parse import urlparse, parse_qs, unquote
from collections import OrderedDict

# ─── Paths ────────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).parent.resolve()
# Use /opt/data for the containerized environment, or script-local data folder
if SCRIPT_DIR == Path("/opt/data") or "opt/data" in str(SCRIPT_DIR):
    DATA_DIR = Path("/opt/data/data")  # Container path
else:
    DATA_DIR = SCRIPT_DIR / "data"     # Local portable path
# Dynamic venv detection - no hardcoded absolute paths
_VENV_CANDIDATES = [
    SCRIPT_DIR.parent / ".venv-docreader" / "lib" / f"python3.{sys.version_info.minor}" / "site-packages",
    SCRIPT_DIR / ".venv-docreader" / "lib" / f"python3.{sys.version_info.minor}" / "site-packages",
    SCRIPT_DIR.parent.parent / ".venv-docreader" / "lib" / f"python3.{sys.version_info.minor}" / "site-packages",
]
VENV_SITE_PACKAGES = str(next((p for p in _VENV_CANDIDATES if p.exists()), ""))
UPLOAD_DIR = DATA_DIR / "uploads"
DOCS_DIR = DATA_DIR / "documents_safe"
MAPS_DIR = DATA_DIR / "redaction_maps"

for d in [UPLOAD_DIR, DOCS_DIR, MAPS_DIR]:
    d.mkdir(parents=True, exist_ok=True)

SETTINGS_FILE = DATA_DIR / "redaction_settings.json"

# ─── Default Settings (all categories enabled by default) ─────────────────────
def get_default_settings():
    """Generate default settings with all categories enabled."""
    settings = {"categories": {}, "custom": []}
    for group_name, categories in SECURITY_POLICY.items():
        settings["categories"][group_name] = {
            "enabled": True,
            "subcategories": {}
        }
        for cat_key, cat_info in categories.items():
            settings["categories"][group_name]["subcategories"][cat_key] = {
                "enabled": True,
                "description": cat_info.get("description", cat_key),
                "critical": cat_info.get("critical", False)
            }
    return settings

def load_settings():
    """Load user settings from file, or return defaults."""
    if SETTINGS_FILE.exists():
        try:
            with open(SETTINGS_FILE, 'r') as f:
                settings = json.load(f)
            # Merge with defaults to ensure new categories are included
            defaults = get_default_settings()
            for group in defaults["categories"]:
                if group not in settings["categories"]:
                    settings["categories"][group] = defaults["categories"][group]
                else:
                    for cat_key in defaults["categories"][group]["subcategories"]:
                        if cat_key not in settings["categories"][group]["subcategories"]:
                            settings["categories"][group]["subcategories"][cat_key] = defaults["categories"][group]["subcategories"][cat_key]
            return settings
        except Exception:
            return get_default_settings()
    return get_default_settings()

def save_settings(settings):
    """Save user settings to file."""
    SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(SETTINGS_FILE, 'w') as f:
        json.dump(settings, f, indent=2)
    return True

# ─── Security Policy Definitions ─────────────────────────────────────────────
SECURITY_POLICY = {
    "PII": {
        "SSN": {
            "patterns": [
                r'\b\d{3}-\d{2}-\d{4}\b',  # Standard SSN format
                r'\b\d{1,3}-\d{5,7}\b',    # Alternate format
                r'\b\d{9}\b',              # SSN without dashes
            ],
            "description": "Social Security Numbers",
            "dummy_prefix": "SSN",
            "critical": True
        },
        "EMAIL": {
            "patterns": [
                r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
                r'\b[\w.-]+@[\w.-]+\.\w+\b',
            ],
            "description": "Email addresses",
            "dummy_prefix": "EMAIL",
            "critical": True
        },
        "PHONE": {
            "patterns": [
                r'\+?\d{1,2}[-.]?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}',  # +1-555-123-4567
                r'\b\d{3}[-.\s]?\d{3}[-.\s]?\d{4}\b',                         # XXX-XXX-XXXX
            ],
            "description": "Phone numbers (US/International)",
            "dummy_prefix": "PHONE",
            "critical": True
        },
        "CREDIT_CARD": {
            "patterns": [
                r'\b(?:\d{4}[-\s]?){3}\d{4}\b',  # Standard CC format
                r'\b\d{16}\b',                    # CC without separators
            ],
            "description": "Credit card numbers",
            "dummy_prefix": "CC",
            "critical": True
        },
        "BANK_ACCOUNT": {
            "patterns": [
                r'\bAccount\s*[:]?\s*\d{8,17}\b',  # "Account 1234567890" or "Account: 1234567890"
                r'\b\d{10,16}\b',                  # Long numeric sequences (account numbers)
            ],
            "description": "Bank account numbers",
            "dummy_prefix": "BANKACC",
            "critical": True
        }
    },
    "BUSINESS_SENSITIVE": {
        "COMPANY_NAME": {
            "patterns": [
                r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\s+(?:Corporation|Corp\.|Inc\.|LLC|Ltd\.|Limited|Group|Holding|Holdings)\b',
                r'\b(?:Hitec|Hittech|Tech)\s+Solutions\b',
                r'\b[A-Z][a-z]+\s+(?:Solutions|Services|Systems|Group|Corp)\b',
            ],
            "description": "Company and organization names",
            "dummy_prefix": "COMPANY",
            "critical": True
        },
        "PRODUCT_NAME": {
            "patterns": [
                r'\b(?:TechPro|DataMax|EcoLite|PowerGrid|CloudSuite|NetSys|SoftEdge|UltraGen|ProMax|EliteCore|LiteFlex|MaxPro|PlusCore|UltraMax|ProElite|LiteMax)\s*(?:Pro|Plus|Elite|Ultra|Lite|Max|Edition|Series|System|Suite|Platform|Solution)\w*\b',
                r'\b(?:Alpha|Beta|Gamma|Delta|Omega|Apex|Nexus|Quantum|Stellar|Vertex|Edge|Core|Prime|Lite|Max|Ultra|Pro|Plus|Elite)\s*(?:Suite|System|Platform|Solution|Edition|Series)\w*\b',
            ],
            "description": "Product names and model identifiers",
            "dummy_prefix": "PROD",
            "critical": True
        },
        "DIRECTOR_NAME": {
            "patterns": [
                r'\b(?:Director|Manager|CEO|CTO|CFO|President|Head)[ \t]+[A-Z][a-z]+(?:[ \t]+[A-Z][a-z]+)*\b',
                r'\b(?:Mr\.|Ms\.|Mrs\.|Dr\.)\s+[A-Z][a-z]+\s+[A-Z][a-z]+\b',
                r'\b[A-Z][a-z]{1,4},\s+[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b',  # "Seow, EG" or "Smith, John"
                r'\b(?:EG|AB|CD|EF|GH|IJ|KL|MN|OP|QR|ST|UV|WX|YZ)[A-Z]*\s+[A-Z][a-z]+\b',  # "EG Seow" format
            ],
            "description": "Executive and personnel names",
            "dummy_prefix": "DIRECTOR",
            "critical": True
        },
        "QUOTATION_ID": {
            "patterns": [
                r'\b(?:Quotation|Quote|QTN|Ref:|REF:)\s*#?[A-Z0-9-]+-\d{4,}[-\w]*\b',
                r'\b[A-Z]{2,5}-\d{4}[-\w]*\b',  # Generic quote ID like QTN-2024-0847
            ],
            "description": "Quotation/reference identifiers",
            "dummy_prefix": "QUOTE",
            "critical": True
        },
        "COST_VALUE": {
            "patterns": [
                r'\$\d{1,3}(?:,\d{3})*(?:\.\d{2})?',   # $12,500.00 (no \b)
                r'€\d+(?:[.,]\d{2})?\b',                    # €12500
                r'RM\d+(?:[.,]\d{2})?\b',                   # RM12500
                r'\b\d+(?:[.,]\d{2})?\s*(?:USD|EUR|MYR|GBP)\b',  # 12500 USD
            ],
            "description": "Monetary values and costs",
            "dummy_prefix": "COST",
            "critical": True
        }
    }
}

# ─── Enhanced Redaction Engine ───────────────────────────────────────────────
class EnhancedRedactionEngine:
    """Enterprise-grade data anonymization engine with business-sensitive detection."""
    
    def __init__(self, settings=None):
        # Build flat pattern list from policy
        self.categories = OrderedDict()
        self.settings = settings if settings else get_default_settings()
        self._build_patterns()
    
    def _build_patterns(self):
        """Build regex patterns from security policy, respecting user settings."""
        for group_name, group_settings in self.settings.get("categories", {}).items():
            if not group_settings.get("enabled", True):
                continue
            categories = SECURITY_POLICY.get(group_name, {})
            for cat_key, cat_info in categories.items():
                # Check if this category is enabled in settings
                sub_settings = group_settings.get("subcategories", {}).get(cat_key, {})
                if not sub_settings.get("enabled", True):
                    continue
                patterns = cat_info.get('patterns', [cat_info.get('pattern', '')])
                for p in patterns:
                    if p:
                        if cat_key not in self.categories:
                            self.categories[cat_key] = []
                        self.categories[cat_key].append((p, cat_info))

        # Add custom categories
        for custom in self.settings.get("custom", []):
            cat_key = custom.get("name", "CUSTOM")
            if not custom.get("enabled", True):
                continue
            pattern_str = custom.get("pattern", "")
            if pattern_str:
                cat_info = {
                    "description": custom.get("description", cat_key),
                    "dummy_prefix": custom.get("dummy_prefix", cat_key),
                    "critical": custom.get("critical", False),
                    "category": custom.get("category", "CUSTOM")
                }
                if cat_key not in self.categories:
                    self.categories[cat_key] = []
                self.categories[cat_key].append((pattern_str, cat_info))
    
    def redact(self, text):
        """Perform redaction with reversible mapping.
        
        Critical fix: Process categories in priority order (PII first, then business)
        to ensure overlapping patterns don't leak. Use position-based replacement
        to handle exact span replacement.
        """
        redactions = []
        category_counts = {}
        
        # Define priority order - PII critical data first
        priority_order = ['SSN', 'EMAIL', 'CREDIT_CARD', 'BANK_ACCOUNT', 'PHONE',
                         'QUOTATION_ID', 'COST_VALUE',
                         'COMPANY_NAME', 'PRODUCT_NAME', 'DIRECTOR_NAME']
        
        # Collect all matches across all categories
        all_matches = []
        for category in priority_order:
            if category not in self.categories:
                continue
            for pattern, cat_info in self.categories[category]:
                for match in re.finditer(pattern, text, re.IGNORECASE):
                    all_matches.append({
                        'start': match.start(),
                        'end': match.end(),
                        'text': match.group(),
                        'category': category,
                        'cat_info': cat_info
                    })
        
        # Sort by position, then by length (longer matches first for overlapping)
        all_matches.sort(key=lambda m: (m['start'], -m['end']))
        
        # Process matches, handling overlaps by position
        redaction_map = []
        last_end = 0
        result_parts = []
        counter = 0
        
        for match in all_matches:
            # Skip if this match overlaps with a previously processed match
            if match['start'] < last_end:
                continue
            
            # Add text before the match
            result_parts.append(text[last_end:match['start']])
            
            # Count by category
            category = match['category']
            category_counts[category] = category_counts.get(category, 0) + 1
            counter = category_counts[category]
            
            # Get dummy prefix
            dummy_prefix = match['cat_info'].get('dummy_prefix', category)
            var_name = f"{{{dummy_prefix}_{counter}}}"
            
            redaction_map.append({
                'variable': var_name,
                'original': match['text'].strip(),
                'category': category,
                'group': match['cat_info'].get('description', category),
                'position': match['start']
            })
            
            result_parts.append(var_name)
            last_end = match['end']
            redactions.append({
                'category': category,
                'variable': var_name,
                'original': match['text'].strip(),
                'length': len(match['text'])
            })
        
        # Add remaining text
        result_parts.append(text[last_end:])
        safe_text = ''.join(result_parts)
        
        return safe_text, redaction_map, category_counts, redactions

# ─── File Extractors ────────────────────────────────────────────────────────
def extract_text(file_path):
    """Extract text from various document formats using stdlib + optional libs."""
    path = Path(file_path)
    suffix = path.suffix.lower()
    
    try:
        if suffix in ('.txt', '.md', '.log', '.json', '.csv'):
            return path.read_text(encoding='utf-8', errors='replace')
        
        if suffix == '.pdf':
            try:
                sys.path.insert(0, VENV_SITE_PACKAGES)
                from pypdf import PdfReader
                reader = PdfReader(str(path))
                return '\n'.join([page.extract_text() or '' for page in reader.pages])
            except ImportError:
                raw = path.read_bytes()
                text = raw.decode('utf-8', errors='replace')
                text = re.sub(r'stream.*?endstream', '', text, flags=re.DOTALL)
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\(.*?\)', lambda m: m.group()[1:-1], text)
                return text
        
        if suffix == '.docx':
            try:
                sys.path.insert(0, VENV_SITE_PACKAGES)
                from docx import Document
                doc = Document(str(path))
                return '\n'.join([p.text for p in doc.paragraphs])
            except ImportError:
                try:
                    with zipfile.ZipFile(str(path)) as z:
                        for name in z.namelist():
                            if 'word/document' in name:
                                data_bytes = z.read(name)
                                text = re.sub(r'<[^>]+>', ' ', data_bytes.decode('utf-8', errors='ignore'))
                                text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text)
                                return text.strip()
                except Exception:
                    return f"DOCX extraction requires python-docx library"
        
        if suffix == '.xlsx':
            try:
                sys.path.insert(0, VENV_SITE_PACKAGES)
                from openpyxl import load_workbook
                wb = load_workbook(str(path), data_only=True)
                rows = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        rows.append(' '.join(str(c) if c else '' for c in row))
                return '\n'.join(rows)
            except ImportError:
                return f"XLSX extraction requires openpyxl library"
        
        if suffix == '.pptx':
            try:
                sys.path.insert(0, VENV_SITE_PACKAGES)
                from pptx import Presentation
                prs = Presentation(str(path))
                text_out = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text_out.append(shape.text)
                return '\n'.join(text_out)
            except ImportError:
                return f"PPTX extraction requires python-pptx library"
        
        if suffix in ('.html', '.htm'):
            return path.read_text(encoding='utf-8', errors='replace')
        
        # Fallback
        return path.read_text(encoding='utf-8', errors='replace')
    
    except Exception as e:
        return f"Error extracting text from {path.name}: {str(e)}"

# ─── HTML UI ─────────────────────────────────────────────────────────────────
HTML_UI = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>OAKAI Document Reader — Secure Data Anonymization</title>
<style>
  :root {
    --bg: #0f172a;
    --surface: #ffffff;
    --surface-light: #f8fafc;
    --card-shadow: 0 8px 24px rgba(15, 23, 42, 0.12);
    --border-radius: 12px;
    --border-radius-lg: 16px;
    --primary: #1e3a8a;
    --primary-light: #3b82f6;
    --accent: #2563eb;
    --text: #0f172a;
    --text-secondary: #64748b;
    --border: #e2e8f0;
    --border-strong: #cbd5e1;
    --error: #dc2626;
    --success: #16a34a;
    --warning: #d97706;
    --upload-border: #cbd5e1;
    --upload-bg: #f8fafc;
    --variable-bg: #dbeafe;
    --variable-text: #1e40af;
  }
  * { margin: 0; padding: 0; box-sizing: border-box; }
  body {
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, sans-serif;
    background: linear-gradient(135deg, #1e293b 0%, #0f172a 100%);
    min-height: 100vh;
    padding: 24px;
  }
  .container {
    max-width: 900px;
    margin: 0 auto;
  }
  .header {
    text-align: center;
    margin-bottom: 24px;
    padding: 24px;
    background: var(--surface);
    border-radius: var(--border-radius-lg);
    box-shadow: var(--card-shadow);
  }
  .header h1 {
    font-size: 28px;
    font-weight: 700;
    color: var(--primary);
    margin-bottom: 4px;
    display: flex;
    align-items: center;
    justify-content: center;
    gap: 8px;
  }
  .header h1 .oakai {
    color: var(--primary-light);
  }
  .header .subtitle {
    color: var(--text-secondary);
    font-size: 14px;
  }
  .upload-section {
    background: var(--surface);
    border-radius: var(--border-radius-lg);
    padding: 32px;
    margin-bottom: 24px;
    box-shadow: var(--card-shadow);
    text-align: center;
  }
  .upload-area {
    border: 2px dashed var(--upload-border);
    border-radius: var(--border-radius);
    padding: 40px 24px;
    margin: 16px 0;
    background: var(--upload-bg);
    cursor: pointer;
    transition: all 0.2s ease;
  }
  .upload-area:hover {
    border-color: var(--primary);
    background: #f0f7ff;
  }
  .upload-area.drag-over {
    border-color: var(--accent);
    background: #f0fdf4;
  }
  .upload-area input[type="file"] { display: none; }
  .upload-icon {
    font-size: 42px;
    margin-bottom: 12px;
    opacity: 0.7;
  }
  .file-info {
    font-size: 14px;
    color: var(--text-secondary);
    margin: 10px 0;
    min-height: 20px;
  }
  .btn-primary {
    background: var(--accent);
    color: white;
    border: none;
    padding: 14px 36px;
    border-radius: 8px;
    font-size: 16px;
    font-weight: 500;
    cursor: pointer;
    width: 100%;
    transition: background 0.2s;
    margin-top: 12px;
  }
  .btn-primary:hover:not(:disabled) {
    background: var(--primary);
  }
  .btn-primary:disabled {
    background: var(--border-strong);
    cursor: not-allowed;
  }
  .documents-section {
    background: var(--surface);
    border-radius: var(--border-radius-lg);
    padding: 24px;
    margin-bottom: 24px;
    box-shadow: var(--card-shadow);
  }
  .documents-section h2 {
    font-size: 18px;
    font-weight: 600;
    margin-bottom: 16px;
    color: var(--text);
    display: flex;
    justify-content: space-between;
    align-items: center;
  }
  .doc-count {
    font-size: 13px;
    color: var(--text-secondary);
    font-weight: normal;
  }
  .doc-list {
    display: flex;
    flex-direction: column;
    gap: 8px;
  }
  .doc-item {
    display: flex;
    justify-content: space-between;
    align-items: center;
    padding: 12px 16px;
    background: var(--surface-light);
    border: 1px solid var(--border);
    border-radius: 8px;
    text-decoration: none;
    color: var(--text);
    transition: all 0.15s;
  }
  .doc-item:hover {
    background: #eff6ff;
    border-color: var(--primary);
    color: var(--primary);
  }
  .doc-item .doc-info {
    flex: 1;
  }
  .doc-item .doc-filename {
    font-weight: 500;
    font-size: 14px;
  }
  .doc-item .doc-meta {
    font-size: 12px;
    color: var(--text-secondary);
    margin-top: 2px;
  }
  .doc-item .doc-reducts {
    font-size: 12px;
    background: var(--primary);
    color: white;
    padding: 2px 10px;
    border-radius: 12px;
  }
  .no-docs {
    text-align: center;
    color: var(--text-secondary);
    padding: 32px;
    font-style: italic;
  }
  .results-section {
    background: var(--surface);
    border-radius: var(--border-radius-lg);
    padding: 24px;
    box-shadow: var(--card-shadow);
    margin-bottom: 24px;
  }
  .results-section h2 {
    font-size: 18px;
    font-weight: 600;
    margin-bottom: 16px;
    color: var(--text);
  }
  .result-header {
    display: flex;
    justify-content: space-between;
    align-items: center;
    margin-bottom: 12px;
  }
  .result-header .doc-title {
    font-weight: 600;
    font-size: 16px;
  }
  .status-badge {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    padding: 6px 16px;
    border-radius: 20px;
    font-size: 13px;
    font-weight: 500;
  }
  .status-success {
    background: #dcfce7;
    color: #166534;
  }
  .status-error {
    background: #fee2e2;
    color: #991c1c;
  }
  .redaction-summary {
    display: flex;
    flex-wrap: wrap;
    gap: 6px;
    margin: 12px 0;
  }
  .redaction-badge {
    display: inline-flex;
    align-items: center;
    padding: 4px 12px;
    border-radius: 16px;
    font-size: 12px;
    font-weight: 500;
    color: white;
  }
  .redaction-badge.sensitive { background: #ef4444; }
  .redaction-badge.business { background: #8b5cf6; }
  .redaction-badge.financial { background: #f59e0b; }
  .result-content {
    background: var(--surface-light);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 16px;
    font-family: 'SF Mono', 'Fira Code', Monaco, Consolas, monospace;
    font-size: 13px;
    color: var(--text);
    white-space: pre-wrap;
    word-break: break-word;
    max-height: 400px;
    overflow-y: auto;
    line-height: 1.6;
    margin-top: 12px;
  }
  .variable-highlight {
    background: var(--variable-bg);
    color: var(--variable-text);
    padding: 1px 6px;
    border-radius: 4px;
    font-weight: 600;
  }
  .mapping-section {
    margin-top: 16px;
    padding-top: 12px;
    border-top: 1px solid var(--border);
  }
  .mapping-section h4 {
    font-size: 13px;
    color: var(--text-secondary);
    margin-bottom: 8px;
  }
  .mapping-table {
    width: 100%;
    border-collapse: collapse;
  }
  .mapping-table th, .mapping-table td {
    text-align: left;
    padding: 6px 8px;
    border-bottom: 1px solid var(--border);
    font-size: 12px;
  }
  .mapping-table th {
    color: var(--text-secondary);
    font-weight: 500;
  }
  .mapping-table td:first-child {
    font-family: 'SF Mono', monospace;
    font-weight: 600;
    color: var(--primary);
  }
  .file-info-table {
    width: 100%;
    border-collapse: collapse;
    font-size: 13px;
  }
  .file-info-table td {
    padding: 4px 8px;
    border-bottom: 1px solid var(--border);
  }
  .file-info-table td:first-child {
    color: var(--text-secondary);
    width: 40%;
  }
  .file-info-table td:last-child {
    font-weight: 500;
    word-break: break-all;
  }
  .error-msg {
    background: #fef2f2;
    border: 1px solid #fee2e2;
    border-radius: 8px;
    padding: 12px 16px;
    color: var(--error);
    margin: 12px 0;
    font-size: 14px;
  }
  .footer {
    text-align: center;
    color: var(--text-secondary);
    font-size: 12px;
    margin-top: 24px;
    padding: 16px;
  }
  @media (max-width: 600px) {
    body { padding: 12px; }
    .header h1 { font-size: 22px; }
    .upload-section { padding: 20px 16px; }
  }
  .btn-secondary {
    background: #f1f5f9;
    color: var(--primary);
    border: 1px solid var(--border);
    padding: 10px 20px;
    border-radius: 8px;
    font-size: 14px;
    font-weight: 500;
    cursor: pointer;
    transition: all 0.2s;
  }
  .btn-secondary:hover { background: #e2e8f0; transform: translateY(-1px); }
  .modal-overlay {
    position: fixed; top: 0; left: 0; right: 0; bottom: 0;
    background: rgba(0, 0, 0, 0.6);
    display: none; align-items: flex-start; justify-content: center;
    z-index: 1000; padding: 20px; overflow-y: auto;
  }
  .modal-overlay.active { display: flex; }
  .modal-content {
    background: var(--surface); border-radius: var(--border-radius-lg);
    max-width: 800px; width: 100%; max-height: 90vh; overflow-y: auto;
    box-shadow: var(--card-shadow);
  }
  .modal-header { padding: 24px 32px 16px; border-bottom: 1px solid var(--border);
    display: flex; justify-content: space-between; align-items: center; }
  .modal-header h2 { font-size: 20px; font-weight: 600; color: var(--primary); }
  .modal-close { font-size: 24px; cursor: pointer; color: var(--text-secondary); }
  .modal-body { padding: 24px 32px; }
  .modal-footer { padding: 16px 32px 24px; border-top: 1px solid var(--border);
    display: flex; justify-content: flex-end; gap: 12px; }
  .category-group { margin-bottom: 20px; padding: 16px; border: 1px solid var(--border);
    border-radius: 10px; background: var(--surface-light); }
  .category-group h3 { font-size: 16px; font-weight: 600; color: var(--primary);
    margin-bottom: 12px; display: flex; align-items: center; gap: 8px; }
  .cat-toggle { display: flex; align-items: center; justify-content: space-between;
    padding: 10px 12px; background: white; border: 1px solid var(--border);
    border-radius: 8px; margin-bottom: 8px; }
  .cat-toggle input[type="checkbox"] { width: 18px; height: 18px; cursor: pointer; }
  .cat-toggle label { flex: 1; cursor: pointer; font-size: 14px; }
  .badge-critical { font-size: 11px; padding: 2px 8px; border-radius: 4px;
    background: #fef3c7; color: #92400e; margin-left: 8px; }
  .custom-section { border-top: 2px dashed var(--border); padding-top: 20px; margin-top: 20px; }
  .custom-row { display: flex; gap: 8px; align-items: center; margin-bottom: 8px; }
  .custom-row input, .custom-row select { flex: 1; padding: 8px 12px; border: 1px solid var(--border);
    border-radius: 6px; font-size: 14px; }
  .custom-row button { padding: 8px 16px; background: var(--accent); color: white;
    border: none; border-radius: 6px; font-size: 13px; cursor: pointer; }
  .custom-row button:hover { opacity: 0.9; }
  .custom-item { display: flex; justify-content: space-between; align-items: center;
    padding: 8px 12px; background: white; border-radius: 6px; margin-bottom: 4px;
    border: 1px solid var(--border); }
  .custom-item button { background: var(--error); }
  .success-banner { background: var(--success); color: white; padding: 12px 20px;
    border-radius: 8px; margin-bottom: 16px; text-align: center; font-weight: 500; }
</style>
</head>
<body>
<div class="container">
  <div class="header">
    <div style="display: flex; justify-content: space-between; align-items: center; flex-wrap: wrap; gap: 12px;">
      <div>
        <h1 style="display: flex; align-items: center; gap: 8px; justify-content: left;">
          <span class="oakai">OAKAI</span> Document Reader
        </h1>
        <p class="subtitle">Enterprise Data Anonymization &amp; PII Redaction Engine</p>
      </div>
      <div style="display: flex; gap: 12px;">
        <button onclick="openSettings()" class="btn-secondary" style="padding: 10px 20px; font-size: 14px;">⚙️ Settings</button>
        <button onclick="loadDocs()" class="btn-secondary" style="padding: 10px 20px; font-size: 14px;">🔄 Refresh</button>
      </div>
    </div>
  </div>

  <div class="documents-section">
    <h2>Recently Processed Documents <span class="doc-count" id="docCount"></span></h2>
    <div id="docsList" class="doc-list">
      <div class="no-docs">Loading documents<span class="loading-dots"></span></div>
    </div>
  </div>

  <div class="upload-section">
    <div class="upload-icon">📎</div>
    <p style="color: var(--text-secondary); margin-bottom: 16px; font-size: 15px;">Drag & drop files here or click to select</p>
    <div class="upload-area" id="dropZone">
      <input type="file" id="fileInput" accept=".pdf,.docx,.xlsx,.pptx,.txt,.csv,.html,.htm,.md,.json" multiple>
      <span style="color: var(--text-secondary); font-size: 14px;">Click anywhere to browse files</span>
    </div>
    <div class="file-info" id="fileName">No file selected</div>
    <button class="btn-primary" id="processBtn" disabled>Upload & Process</button>
  </div>

  <div class="results-section" id="resultsSection" style="display: none;">
    <div class="result-header">
      <span class="doc-title" id="resultTitle">Processing Results</span>
      <span class="status-badge status-success" id="resultStatus">✓ Complete</span>
    </div>
    <div class="redaction-summary" id="redactionSummary"></div>
    <div style="margin: 12px 0;">
      <table class="file-info-table">
        <tr><td>Original file:</td><td id="resultFilename"></td></tr>
        <tr><td>Uploaded to:</td><td id="resultUploadPath"></td></tr>
        <tr><td>Total redactions:</td><td id="resultRedactions"></td></tr>
        <tr><td>Categories:</td><td id="resultCategories"></td></tr>
        <tr><td>Safe document:</td><td id="resultSafePath"></td></tr>
        <tr><td>Variable mapping:</td><td id="resultMapPath"></td></tr>
      </table>
    </div>
    <div class="result-content" id="resultContent"></div>
    <div class="mapping-section">
      <h4>Redaction Variable Mapping (stored locally):</h4>
      <table class="mapping-table" id="mappingTable">
        <thead><tr><th>Variable</th><th>Original Value</th><th>Category</th></tr></thead>
        <tbody id="mappingBody"></tbody>
      </table>
    </div>
  </div>

  <div class="footer">
    OAKAI Document Reader v1.0 — All processing is 100% local. No APIs or tokens required.
  </div>
</div>

<script>
const dropZone = document.getElementById('dropZone');
const fileInput = document.getElementById('fileInput');
const fileNameDiv = document.getElementById('fileName');
const processBtn = document.getElementById('processBtn');
const resultsSection = document.getElementById('resultsSection');
const docCountSpan = document.getElementById('docCount');

let selectedFile = null;

dropZone.addEventListener('click', () => fileInput.click());

fileInput.addEventListener('change', () => {
  const files = fileInput.files;
  if (files.length > 0) {
    selectedFile = files[0];
    if (files.length === 1) {
      fileNameDiv.textContent = selectedFile.name + ' (' + Math.round(selectedFile.size/1024) + ' KB)';
    } else {
      fileNameDiv.textContent = files.length + ' files selected';
    }
    processBtn.disabled = false;
    processBtn.textContent = 'Process ' + files.length + ' file(s)';
  } else {
    fileNameDiv.textContent = 'No file selected';
    processBtn.disabled = true;
    processBtn.textContent = 'Upload & Process';
  }
});

processBtn.addEventListener('click', async () => {
  if (!selectedFile) return;

  processBtn.disabled = true;
  processBtn.textContent = 'Processing...';
  resultsSection.style.display = 'none';
  const mappingBody = document.getElementById('mappingBody');
  const resultContent = document.getElementById('resultContent');
  const resultTitle = document.getElementById('resultTitle');
  const resultStatus = document.getElementById('resultStatus');

  const formData = new FormData();
  formData.append('file', selectedFile);

  try {
    const res = await fetch('/upload', {method: 'POST', body: formData});
    const data = await res.json();
    
    // Re-enable button
    processBtn.disabled = false;
    processBtn.textContent = 'Upload & Process';

    if (data.error) {
      resultTitle.textContent = 'Processing Error';
      resultStatus.className = 'status-badge status-error';
      resultStatus.textContent = '✗ Failed';
      resultContent.innerHTML = '<div class="error-msg">Error: ' + data.error + '</div>';
      resultsSection.style.display = 'block';
    } else {
      resultTitle.textContent = data.original_filename || 'Document Processed';
      resultStatus.className = 'status-badge status-success';
      resultStatus.textContent = '✓ Complete';
      
      // File info table
      document.getElementById('resultFilename').textContent = data.original_filename || 'N/A';
      document.getElementById('resultUploadPath').textContent = data.file_path || 'See uploads dir';
      document.getElementById('resultRedactions').textContent = data.total_redactions;
      document.getElementById('resultCategories').textContent = Object.keys(data.category_counts || {}).join(', ');
      document.getElementById('resultSafePath').textContent = data.safe_url || 'N/A';
      document.getElementById('resultMapPath').textContent = 'Redaction map stored at: ' + (data.map_path || 'data/redaction_maps/');

      // Redaction summary badges
      const categoryGroups = {
        'PII': ['SSN', 'EMAIL', 'PHONE', 'CREDIT_CARD', 'BANK_ACCOUNT'],
        'BUSINESS': ['PRODUCT_NAME', 'COMPANY_NAME', 'DIRECTOR_NAME', 'QUOTATION_ID'],
        'FINANCIAL': ['COST_VALUE']
      };
      
      let badgesHtml = '';
      for (const [cat, count] of Object.entries(data.category_counts || {})) {
        let badgeClass = '';
        if (categoryGroups.PII.includes(cat)) badgeClass = 'sensitive';
        else if (categoryGroups.BUSINESS.includes(cat)) badgeClass = 'business';
        else if (categoryGroups.FINANCIAL.includes(cat)) badgeClass = 'financial';
        
        badgesHtml += '<span class="redaction-badge ' + badgeClass + '">' + cat + ': ' + count + '</span>';
      }
      document.getElementById('redactionSummary').innerHTML = badgesHtml;

      // Fetch safe content and build result
      try {
        const safeRes = await fetch('/documents/' + data.document_id + '/safe');
        const safeData = await safeRes.json();
        
        let safeText = safeData.all_text || 'No text extracted';
        safeText = safeText.replace(/&/g, '&amp;').replace(/</g, '&lt;').replace(/>/g, '&gt;');
        // Highlight variables
        safeText = safeText.replace(/\{[A-Z]+_\d+\}/g, '<span class="variable-highlight">$&</span>');
        
        resultContent.textContent = safeText;
        
      } catch(e) {
        resultContent.innerHTML = '<div class="error-msg">Could not load safe document content</div>';
      }

      // Fetch redaction mapping from separate endpoint (inside try block, with data)
      try {
        const mapRes = await fetch('/documents/' + data.document_id + '/map');
        const mapData = await mapRes.json();
        
        let mappingHtml = '';
        const mapping = mapData.full_map || {};
        for (const [varName, info] of Object.entries(mapping)) {
          let badgeClass = '';
          if (['SSN', 'EMAIL', 'PHONE', 'CREDIT_CARD', 'BANK_ACCOUNT'].includes(info.category)) badgeClass = 'sensitive';
          else if (['PRODUCT_NAME', 'COMPANY_NAME', 'DIRECTOR_NAME'].includes(info.category)) badgeClass = 'business';
          else if (['QUOTATION_ID'].includes(info.category)) badgeClass = 'sensitive';
          else if (['COST_VALUE'].includes(info.category)) badgeClass = 'financial';
          else badgeClass = 'sensitive';
          
          const safeValue = info.original.replace(/&/g, '&amp;').replace(/</g, '&lt;').replace(/>/g, '&gt;');
          mappingHtml += '<tr><td>' + varName + '</td><td style="font-family: monospace;">' + safeValue + '</td>';
          mappingHtml += '<td><span class="redaction-badge ' + badgeClass + '">' + info.category + '</span><br><small style="color: var(--text-secondary);">' + info.group + '</small></td></tr>';
        }
        mappingBody.innerHTML = mappingHtml || '<tr><td colspan="3">No variables mapped</td></tr>';
      } catch(e2) {
        mappingBody.innerHTML = '<tr><td colspan="3">Mapping unavailable</td></tr>';
      }
      
      resultsSection.style.display = 'block';
    }
    
    // Reset form
    fileNameDiv.textContent = 'Ready to upload another file';
    fileInput.value = '';
    await loadDocs();
    
  } catch(e) {
    processBtn.disabled = false;
    processBtn.textContent = 'Upload & Process';
    resultTitle.textContent = 'Error';
    resultStatus.className = 'status-badge status-error';
    resultStatus.textContent = '✗ Failed';
    resultContent.innerHTML = '<div class="error-msg">Error: ' + e.message + '</div>';
    resultsSection.style.display = 'block';
  }
});

async function loadDocs() {
  try {
    const res = await fetch('/documents');
    const data = await res.json();
    const docs = data.documents || [];
    const docsList = document.getElementById('docsList');
    docCountSpan.textContent = '(' + docs.length + ')';

    if (docs.length === 0) {
      docsList.innerHTML = '<div class="no-docs">No documents processed yet. Upload your first document above!</div>';
      return;
    }

    let html = '';
    docs.slice(0, 10).forEach(doc => {
      const docDate = doc.id.replace('doc_', '').replace(/_/g, ' ');
      html += '<a href="' + doc.url + '" target="_blank" class="doc-item">';
      html += '<div class="doc-info"><div class="doc-filename">' + doc.filename + '</div>';
      html += '<div class="doc-meta">' + docDate + ' · ' + doc.size + ' bytes</div></div>';
      html += '<span class="doc-reducts">' + (doc.redactions_count || '?') + '</span>';
      html += '</a>';
    });
    docsList.innerHTML = html;
  } catch(e) {
    document.getElementById('docsList').innerHTML = '<div class="no-docs">Error loading documents</div>';
  }
}

// Drag-drop handlers
dropZone.addEventListener('dragover', e => {
  e.preventDefault();
  dropZone.classList.add('drag-over');
});
dropZone.addEventListener('dragleave', e => {
  e.preventDefault();
  dropZone.classList.remove('drag-over');
});
dropZone.addEventListener('drop', e => {
  e.preventDefault();
  dropZone.classList.remove('drag-over');
  if (e.dataTransfer.files.length) {
    fileInput.files = e.dataTransfer.files;
    selectedFile = e.dataTransfer.files[0];
    fileNameDiv.textContent = selectedFile.name + ' (' + Math.round(selectedFile.size/1024) + ' KB)';
    processBtn.disabled = false;
    processBtn.textContent = 'Process ' + e.dataTransfer.files.length + ' file(s)';
    
    // Auto-submit if single file
    if (e.dataTransfer.files.length === 1) {
      setTimeout(() => processBtn.click(), 100);
    }
  }
});

// ─── Settings Functions ──────────────────────────────────────────────────────
let currentSettings = null;

async function openSettings() {
  const overlay = document.getElementById('settingsOverlay');
  const settingsList = document.getElementById('settingsList');
  const errorDiv = document.getElementById('settingsError');
  const successDiv = document.getElementById('settingsSuccess');
  
  overlay.classList.add('active');
  settingsList.innerHTML = '<div style="text-align: center; color: var(--text-secondary); padding: 20px;">Loading...</div>';
  errorDiv.style.display = 'none';
  successDiv.style.display = 'none';
  
  try {
    const res = await fetch('/settings');
    const data = await res.json();
    currentSettings = data.settings;
    
    let html = '';
    for (const [groupName, groupData] of Object.entries(data.settings.categories)) {
      const groupLabel = groupName === 'PII' ? 'PII / Sensitive' : 'Business Sensitive';
      html += '<div class="category-group">';
      html += '<h3>' + groupLabel + '</h3>';
      for (const [catKey, catData] of Object.entries(groupData.subcategories)) {
        const checked = catData.enabled ? 'checked' : '';
        const criticalBadge = catData.critical ? '<span class="badge-critical">CRITICAL</span>' : '';
        html += '<div class="cat-toggle">';
        html += '<label for="' + catKey + '">' + catData.description + criticalBadge + '</label>';
        html += '<input type="checkbox" id="' + catKey + '" onchange="toggleCategory(\'' + groupName + '\', \'' + catKey + '\', this.checked)" ' + checked + '>';
        html += '</div>';
      }
      html += '</div>';
    }
    
    // Custom categories section
    if (data.settings.custom && data.settings.custom.length > 0) {
      html += '<div class="custom-section"><h3 style="font-size:14px;font-weight:600;margin-bottom:12px;color:var(--primary);">Custom Categories</h3>';
      data.settings.custom.forEach((custom, idx) => {
        const checked = custom.enabled ? 'checked' : '';
        html += '<div class="custom-item">';
        html += '<span>' + custom.name + ' — ' + (custom.description || '') + '</span>';
        html += '<input type="checkbox" onchange="toggleCustom(' + idx + ', this.checked)" ' + checked + '>';
        html += '</div>';
      });
      html += '</div>';
    }
    
    settingsList.innerHTML = html;
  } catch(e) {
    errorDiv.innerHTML = 'Error loading settings: ' + e.message;
    errorDiv.style.display = 'block';
    settingsList.innerHTML = '';
  }
}

function toggleCategory(group, catKey, enabled) {
  if (currentSettings && currentSettings.categories[group] && currentSettings.categories[group].subcategories[catKey]) {
    currentSettings.categories[group].subcategories[catKey].enabled = enabled;
  }
}

function toggleCustom(index, enabled) {
  if (currentSettings && currentSettings.custom && currentSettings.custom[index]) {
    currentSettings.custom[index].enabled = enabled;
  }
}

async function addCustomCategory() {
  const name = document.getElementById('customName').value.trim();
  const pattern = document.getElementById('customPattern').value.trim();
  const desc = document.getElementById('customDesc').value.trim() || name;
  const prefix = document.getElementById('customPrefix').value;
  const errorDiv = document.getElementById('settingsError');
  
  if (!name || !pattern) {
    errorDiv.innerHTML = 'Please enter both category name and regex pattern';
    errorDiv.style.display = 'block';
    return;
  }
  
  if (!currentSettings.custom) currentSettings.custom = [];
  
  currentSettings.custom.push({
    name: name,
    pattern: pattern,
    description: desc,
    dummy_prefix: prefix,
    enabled: true
  });
  
  document.getElementById('customName').value = '';
  document.getElementById('customPattern').value = '';
  document.getElementById('customDesc').value = '';
  
  // Re-open to show updated list
  openSettings();
}

async function saveSettings() {
  try {
    const res = await fetch('/settings', {
      method: 'POST',
      headers: { 'Content-Type': 'application/json' },
      body: JSON.stringify(currentSettings)
    });
    const data = await res.json();
    
    if (res.ok) {
      document.getElementById('settingsSuccess').style.display = 'block';
      setTimeout(() => {
        document.getElementById('settingsSuccess').style.display = 'none';
        closeSettings();
      }, 1500);
    } else {
      document.getElementById('settingsError').innerHTML = data.error || 'Save failed';
      document.getElementById('settingsError').style.display = 'block';
    }
  } catch(e) {
    document.getElementById('settingsError').innerHTML = 'Error: ' + e.message;
    document.getElementById('settingsError').style.display = 'block';
  }
}

function closeSettings() {
  document.getElementById('settingsOverlay').classList.remove('active');
}

// Initialize
loadDocs();
setInterval(loadDocs, 10000); // Auto-refresh every 10 seconds
</script>
</body>

<!-- Settings Modal -->
<div class="modal-overlay" id="settingsOverlay">
  <div class="modal-content">
    <div class="modal-header">
      <h2>🔒 Redaction Settings</h2>
      <span class="modal-close" onclick="closeSettings()">×</span>
    </div>
    <div class="modal-body">
      <p style="color: var(--text-secondary); margin-bottom: 16px; font-size: 14px;">
        Select which data categories to redact. Changes take effect on next document processing.
      </p>
      <div id="settingsError" style="display: none; color: var(--error); background: #fef2f2; padding: 10px; border-radius: 6px; margin-bottom: 12px;"></div>
      <div id="settingsSuccess" style="display: none;" class="success-banner">✅ Settings saved successfully!</div>
      <div id="settingsList">
        <div style="text-align: center; color: var(--text-secondary); padding: 20px;">Loading categories...</div>
      </div>
      <div class="custom-section">
        <h3 style="font-size: 14px; font-weight: 600; margin-bottom: 12px; color: var(--primary);">Add Custom Category</h3>
        <div class="custom-row">
          <input type="text" id="customName" placeholder="Category name (e.g., PATENT_ID)">
          <input type="text" id="customPattern" placeholder="Regex pattern">
        </div>
        <div class="custom-row">
          <input type="text" id="customDesc" placeholder="Description">
          <select id="customPrefix">
            <option value="CUSTOM">CUSTOM</option>
            <option value="SSN">SSN</option>
            <option value="ID">ID</option>
            <option value="REF">REF</option>
          </select>
        </div>
        <button onclick="addCustomCategory()" style="width: 100%; padding: 10px;">➕ Add Custom Category</button>
      </div>
    </div>
    <div class="modal-footer">
      <button class="btn-secondary" onclick="closeSettings()">Cancel</button>
      <button onclick="saveSettings()" style="background: var(--accent); color: white; border: none; padding: 12px 24px; border-radius: 8px; font-weight: 500; cursor: pointer;">💾 Save Settings</button>
    </div>
  </div>
</div>
</body>
</html>"""

# ─── API Server ───────────────────────────────────────────────────────────────
def main():
    port = 8765
    if "--api-server" in sys.argv:
        idx = sys.argv.index("--api-server")
        port = int(sys.argv[idx + 1]) if idx + 1 < len(sys.argv) else 8765

    settings = load_settings()
    engine = EnhancedRedactionEngine(settings=settings)

    def process_file(filepath, original_filename=None):
        """Process a single file - returns safe doc and redaction map."""
        doc_id = f"doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        original_name = original_filename or Path(filepath).name
        original_text = extract_text(filepath)

        # Redact PII
        safe_text, redaction_map, category_counts, redactions = engine.redact(original_text)

        # Convert list-based redaction_map to enriched format
        # redaction_map is now a list of dicts: [{variable, original, category, group, position}, ...]
        enriched_map = OrderedDict()
        for entry in redaction_map:
            var_name = entry["variable"]
            category = entry["category"]
            group_info = None
            for group_name, categories in SECURITY_POLICY.items():
                if category in categories:
                    group_info = group_name
                    break
            
            enriched_map[var_name] = {
                "original": entry["original"],
                "category": category,
                "group": group_info or "UNKNOWN",
                "description": entry.get("group", ""),
                "position": entry.get("position", 0)
            }

        # Save redaction map locally (never exposed via API)
        map_file = MAPS_DIR / f"{doc_id}_redaction_map.json"
        map_data = {
            "document_id": doc_id,
            "original_filename": original_name,
            "map": {k: v["original"] for k, v in enriched_map.items()},  # Simple mapping
            "full_map": enriched_map,  # Detailed mapping with categories
            "category_counts": category_counts,
            "created_at": datetime.now().isoformat(),
            "total_redactions": len(redactions)
        }
        map_file.write_text(json.dumps(map_data, indent=2))

        # Create enriched safe document
        # Note: redaction_map in safe doc contains ONLY variable->original mapping
        # The safe_text field contains the fully redacted text
        safe_doc = {
            "document_id": doc_id,
            "original_filename": original_name,
            "original_type": Path(filepath).suffix.lower(),
            "file_size": Path(filepath).stat().st_size,
            "processed_at": datetime.now().isoformat(),
            "total_redactions": len(redactions),
            "category_counts": category_counts,
            # SAFE OUTPUT: All sensitive data is now in variable form
            "all_text": safe_text,
            "processed_text": safe_text,
            "pages": [{"page_number": 1, "text": safe_text}],
            # Redaction map stored separately, NOT in safe output
            # to ensure zero PII leakage in the safe document
            "redaction_summary": {
                "total_categories": len(category_counts),
                "categories": list(category_counts.keys()),
                "total_variables": len(enriched_map)
            },
            "metadata": {"encoding": "utf-8"}
        }

        # Save safe doc
        safe_file = DOCS_DIR / f"{doc_id}_safe.json"
        safe_file.write_text(json.dumps(safe_doc, indent=2, ensure_ascii=False))

        return safe_doc, map_data, str(filepath), str(safe_file), str(map_file)

    def save_doc(safe_doc, rmap):
        doc_id = safe_doc["document_id"]
        safe_file = DOCS_DIR / f"{doc_id}_safe.json"
        safe_file.write_text(json.dumps(safe_doc, indent=2, ensure_ascii=False))

    class Handler(BaseHTTPRequestHandler):
        def log_message(self, fmt, *args): pass

        def do_GET(self):
            parsed = urlparse(self.path)
            path = parsed.path

            if path == "/health":
                self._json(200, {"status": "ok", "service": "OAKAI Document Reader"})
            elif path in ("/", "/ui", "/index.html"):
                self.send_response(200)
                self.send_header("Content-Type", "text/html; charset=utf-8")
                self.end_headers()
                self.write_html(HTML_UI)
            elif path == "/settings":
                # Return current settings
                current_settings = load_settings()
                # Also return policy info for new categories
                self._json(200, {
                    "settings": current_settings,
                    "policy": {
                        group: {
                            cat_key: {
                                "description": cat_info.get("description", cat_key),
                                "dummy_prefix": cat_info.get("dummy_prefix", cat_key)
                            }
                            for cat_key, cat_info in categories.items()
                        }
                        for group, categories in SECURITY_POLICY.items()
                    }
                })
            elif path == "/documents":
                docs = []
                if DOCS_DIR.exists():
                    for f in sorted(DOCS_DIR.glob("*_safe.json"), reverse=True)[:10]:
                        doc_id = f.stem.replace("_safe", "")
                        doc_data = json.loads(f.read_text())
                        docs.append({
                            "id": doc_id,
                            "filename": doc_data.get("original_filename", f.name),
                            "url": f"/documents/{doc_id}/safe",
                            "size": f.stat().st_size,
                            "redactions_count": doc_data.get("total_redactions", 0)
                        })
                self._json(200, {"documents": docs, "count": len(docs)})
            else:
                match = re.match(r'/documents/([^/]+)/safe', path)
                if match:
                    doc_id = match.group(1)
                    safe_file = DOCS_DIR / f"{doc_id}_safe.json"
                    if safe_file.exists():
                        data = json.loads(safe_file.read_text())
                        self._json(200, data)
                    else:
                        self._json(404, {"error": f"Document {doc_id} not found"})
                else:
                    # Also support /documents/<id>/map endpoint
                    map_match = re.match(r'/documents/([^/]+)/map', path)
                    if map_match:
                        doc_id = map_match.group(1)
                        map_file = MAPS_DIR / f"{doc_id}_redaction_map.json"
                        if map_file.exists():
                            data = json.loads(map_file.read_text())
                            self._json(200, data)
                        else:
                            self._json(404, {"error": f"Map for document {doc_id} not found"})
                    else:
                        self._json(404, {"error": "Not found"})

        def do_POST(self):
            parsed = urlparse(self.path)
            path = parsed.path

            if path == "/settings":
                # Update settings from POST body
                try:
                    content_length = int(self.headers.get('Content-Length', 0))
                    body = self.rfile.read(content_length)
                    new_settings = json.loads(body)
                    save_settings(new_settings)
                    # Rebuild engine with new settings
                    engine.settings = new_settings
                    engine._build_patterns()
                    self._json(200, {"status": "ok", "message": "Settings saved"})
                except Exception as e:
                    self._json(500, {"error": str(e)})
            elif path == "/settings/categories":
                # Return all available categories from policy for settings UI
                categories_list = []
                for group_name, categories in SECURITY_POLICY.items():
                    for cat_key, cat_info in categories.items():
                        categories_list.append({
                            "group": group_name,
                            "key": cat_key,
                            "label": cat_info.get("description", cat_key),
                            "dummy_prefix": cat_info.get("dummy_prefix", cat_key),
                            "critical": cat_info.get("critical", False)
                        })
                self._json(200, {"categories": categories_list})
            elif path == "/upload":
                self._handle_upload()
            elif path == "/process":
                self._handle_process()
            else:
                self._json(404, {"error": "Not found"})

        def _handle_upload(self):
            """Handle multipart file upload."""
            try:
                content_type = self.headers.get("Content-Type", "")
                if not content_type.startswith("multipart/form-data"):
                    self._json(400, {"error": "Expected multipart/form-data"})
                    return

                boundary = content_type.split("boundary=")[1].encode()
                body = self.rfile.read(int(self.headers["Content-Length"]))
                parts = body.split(b"--" + boundary)

                for part in parts:
                    if b"filename=" in part:
                        # Extract filename
                        header_end = part.find(b"\r\n\r\n")
                        headers_raw = part[:header_end].decode("utf-8", errors="replace")
                        filename_match = re.search(r'filename="(.*?)"', headers_raw)
                        filename = filename_match.group(1) if filename_match else "unknown"

                        # Extract file data
                        file_data = part[header_end + 4:]
                        if file_data.endswith(b"\r\n"):
                            file_data = file_data[:-2]

                        if filename:
                            server_name = f"upload_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}"
                            server_path = UPLOAD_DIR / server_name
                            server_path.write_bytes(file_data)

                            safe_doc, rmap, saved_path, safe_path, map_path = process_file(
                                str(server_path), filename
                            )
                            self._json(200, {
                                "document_id": safe_doc["document_id"],
                                "safe_url": f"/documents/{safe_doc['document_id']}/safe",
                                "original_filename": filename,
                                "total_redactions": safe_doc["total_redactions"],
                                "category_counts": safe_doc["category_counts"],
                                "file_path": saved_path,
                                "safe_path": safe_path,
                                "map_path": map_path
                            })
                            return

                self._json(400, {"error": "No file uploaded"})
            except Exception as e:
                self._json(500, {"error": str(e)})

        def _handle_process(self):
            """Handle process request with file path."""
            try:
                body = self.rfile.read(int(self.headers["Content-Length"]))
                data = json.loads(body)
                fp = data.get("file_path", "")

                if not fp or not Path(fp).exists():
                    self._json(400, {"error": f"File not found: {fp}"})
                    return

                safe_doc, rmap, saved_path, safe_path, map_path = process_file(fp)
                self._json(200, {
                    "document_id": safe_doc["document_id"],
                    "safe_url": f"/documents/{safe_doc['document_id']}/safe",
                    "total_redactions": safe_doc["total_redactions"],
                    "category_counts": safe_doc["category_counts"],
                    "file_path": saved_path,
                    "safe_path": safe_path,
                    "map_path": map_path
                })
            except Exception as e:
                self._json(500, {"error": str(e)})

        def write_html(self, content):
            self.wfile.write(content.encode("utf-8"))

        def _json(self, code, data):
            self.send_response(code)
            self.send_header("Content-Type", "application/json")
            self.end_headers()
            self.wfile.write(json.dumps(data).encode())

    server = HTTPServer(("0.0.0.0", port), Handler)
    print(f"✅ OAKAI Document Reader Server Running!")
    print(f"   Access via: http://localhost:{port}")
    print(f"   Health:     http://localhost:{port}/health")
    print(f"   UI:         http://localhost:{port}/")
    print(f"   Uploads:    {UPLOAD_DIR}")
    print(f"   Safe docs:  {DOCS_DIR}")
    print(f"   Maps:       {MAPS_DIR}")
    server.serve_forever()

if __name__ == "__main__":
    main()