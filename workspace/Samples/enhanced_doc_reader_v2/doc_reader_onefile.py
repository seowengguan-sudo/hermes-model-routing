#!/usr/bin/env python3
"""
OAKAI Document Reader - Enterprise Data Anonymization Engine
=============================================================
Self-contained document reader that performs enterprise-grade data 
anonymization/pseudonymization locally. No LLM/token/API calls.
All processing is 100% local with reversible variable mapping.

Just run: python3 doc_reader_onefile.py
Then open: http://localhost:8765

Supports: PDF, Word, Excel, PPTX, TXT, CSV, HTML, JSON, MD.
"""
import json, os, re, sys, io, zipfile, hashlib, time, subprocess, threading, random
from pathlib import Path
from datetime import datetime
from http.server import HTTPServer, BaseHTTPRequestHandler
from urllib.parse import urlparse, parse_qs, unquote
from collections import OrderedDict

# ─── Paths ────────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).parent.resolve()
# Use /opt/data for the containerized environment, or script-local data folder
# Self-contained project detection: if running from projects/doc_reader/
if "projects/doc_reader" in str(SCRIPT_DIR) or (SCRIPT_DIR.name == "doc_reader" and "projects" in str(SCRIPT_DIR.parent)):
    DATA_DIR = SCRIPT_DIR / "data"     # Project-local portable path
elif SCRIPT_DIR == Path("/opt/data") or "opt/data" in str(SCRIPT_DIR):
    DATA_DIR = Path("/opt/data/data")  # Container path
else:
    DATA_DIR = SCRIPT_DIR / "data"     # Local portable path
# Dynamic venv detection - no hardcoded absolute paths
_VENV_CANDIDATES = [
    SCRIPT_DIR.parent / ".venv-docreader" / "lib" / f"python3.{sys.version_info.minor}" / "site-packages",
    SCRIPT_DIR / ".venv-docreader" / "lib" / f"python3.{sys.version_info.minor}" / "site-packages",
    SCRIPT_DIR.parent.parent / ".venv-docreader" / "lib" / f"python3.{sys.version_info.minor}" / "site-packages",
]
VENV_SITE_PACKAGES = str(next((p for p in _VENV_CANDIDATES if p.exists()), ""))
UPLOAD_DIR = DATA_DIR / "uploads"
DOCS_DIR = DATA_DIR / "documents_safe"
MAPS_DIR = DATA_DIR / "redaction_maps"

for d in [UPLOAD_DIR, DOCS_DIR, MAPS_DIR]:
    d.mkdir(parents=True, exist_ok=True)

SETTINGS_FILE = DATA_DIR / "redaction_settings.json"

# ─── Default Settings (all categories enabled by default) ─────────────────────
def get_default_settings():
    """Generate default settings with all categories enabled."""
    settings = {"categories": {}, "custom": [], "redaction_style": "token"}
    for group_name, categories in SECURITY_POLICY.items():
        settings["categories"][group_name] = {
            "enabled": True,
            "subcategories": {}
        }
        for cat_key, cat_info in categories.items():
            settings["categories"][group_name]["subcategories"][cat_key] = {
                "enabled": True,
                "description": cat_info.get("description", cat_key),
                "critical": cat_info.get("critical", False)
            }
    return settings

def load_settings():
    """Load user settings from file, or return defaults."""
    if SETTINGS_FILE.exists():
        try:
            with open(SETTINGS_FILE, 'r') as f:
                settings = json.load(f)
            # Merge with defaults to ensure new categories are included
            defaults = get_default_settings()
            for group in defaults["categories"]:
                if group not in settings["categories"]:
                    settings["categories"][group] = defaults["categories"][group]
                else:
                    for cat_key in defaults["categories"][group]["subcategories"]:
                        if cat_key not in settings["categories"][group]["subcategories"]:
                            settings["categories"][group]["subcategories"][cat_key] = defaults["categories"][group]["subcategories"][cat_key]
            # Ensure the redaction style key exists (default: token / opaque)
            if "redaction_style" not in settings or settings.get("redaction_style") not in ("token", "smart"):
                settings["redaction_style"] = "token"
            return settings
        except Exception:
            return get_default_settings()
    return get_default_settings()

def save_settings(settings):
    """Save user settings to file."""
    SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(SETTINGS_FILE, 'w') as f:
        json.dump(settings, f, indent=2)
    return True

# ─── Security Policy Definitions ─────────────────────────────────────────────
# PATTERNS USE CAPTURE GROUP 1 for the sensitive VALUE so that, in "Smart Dummy"
# mode, only the value is replaced with a realistic dummy while labels/structure
# (e.g. "account number:", "Name:", "IBAN:") are preserved. In "Token" mode the
# whole match is replaced by a {PREFIX_n} placeholder (original behaviour).
SECURITY_POLICY = {
    "PII": {
        "SSN": {
            "patterns": [
                r'(?P<val>\b\d{3}-\d{2}-\d{4}\b)',
                r'(?P<val>\b\d{9}\b)',
            ],
            "description": "Social Security Numbers",
            "dummy_prefix": "SSN",
            "critical": True
        },
        "EMAIL": {
            "patterns": [ r'(?P<val>[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,})' ],
            "description": "Email addresses",
            "dummy_prefix": "EMAIL",
            "critical": True
        },
        "PHONE": {
            "patterns": [
                r'(?P<val>\+\d{1,3}[\s.-]?\d{2,4}[\s.-]?\d{3,4}[\s.-]?\d{3,4})',
                r'(?P<val>\b0\d{2,3}[\s.-]?\d{3,4}[\s.-]?\d{3,4}\b)',
                r'(?P<val>\b\d{3}[\s.-]?\d{3}[\s.-]?\d{4}\b)',
            ],
            "description": "Phone numbers (local & international)",
            "dummy_prefix": "PHONE",
            "critical": True
        },
        "CREDIT_CARD": {
            "patterns": [ r'(?P<val>\b(?:\d{4}[\s-]?){3}\d{4}\b)' ],
            "description": "Credit card numbers",
            "dummy_prefix": "CC",
            "critical": True
        },
        "BANK_ACCOUNT": {
            "patterns": [
                r'(?P<label>(?:account\s*(?:number|no|#)?|acct|mykad|ic(?:\s*/\s*mykad)?)\s*[:=#]?\s*)(?P<val>\d{6,20})',
            ],
            "description": "Bank account numbers",
            "dummy_prefix": "BANKACC",
            "critical": True
        }
    },
    "BUSINESS_SENSITIVE": {
        "COMPANY_NAME": {
            "patterns": [
                r'(?P<val>\b[A-Z][a-zA-Z]+(?:\s+(?:Corporation|Corp\.?|Inc\.?|LLC|Ltd\.?|Limited|Group|Holdings?|Sdn Bhd|Berhad))\b)',
                r'(?P<val>\b[A-Z][a-z]+(?:\s+(?:Solutions|Services|Systems|Technologies|Industries|Enterprise))\b)',
            ],
            "description": "Company and organization names",
            "dummy_prefix": "COMPANY",
            "critical": True
        },
        "PRODUCT_NAME": {
            "patterns": [
                # Brand-prefixed product names: "<Brand> <Model>" where the brand is a
                # recognised vendor (case-insensitive) and the model contains at least
                # one uppercase letter or digit (case-explicit, so lowercase-only words
                # like "the"/"fruit" fail), and is NOT a company suffix (Inc/Corp/Ltd/…).
                r'(?P<val>\b(?:Samsung|Apple|Dell|HP|Lenovo|Cisco|IBM|Microsoft|Google|Amazon|Huawei|Xiaomi|Sony|LG|Asus|Acer|Intel|AMD|NVIDIA|Oracle|SAP|Salesforce|Tableau|Snowflake|Adobe|VMware|Tesla|Toyota|Honda|BMW|Siemens|Schneider|Bosch|Honeywell|ABB|Rockwell|Omron|Mitsubishi|Yaskawa|Panasonic|Canon|Epson|Brother|Seagate|Western\s*Digital|Kingston|Crucial|Logitech|TP[- ]?Link|D[- ]?Link|Netgear|Fortinet|Palo\s*Alto|Check\s*Point|Juniper|Aruba|Ubiquiti|Zebra|Eaton|APC|Emerson)\b\s+(?!Inc\b|Corp\b|Ltd\b|LLC\b|Group\b|Holdings?\b|Sdn\b|Berhad\b|the\b|and\b|of\b|for\b|with\b|to\b|a\b|an\b|is\b|are\b|was\b|by\b|from\b|at\b|on\b|in\b)(?-i:[A-Za-z0-9]*[A-Z0-9][A-Za-z0-9]*(?:\s+[A-Za-z0-9]+){0,3}))',
                # Model-number SKUs: 2-4 letter prefix + digits (optionally hyphen-
                # joined, or space-joined only when the prefix is 3-4 letters to avoid
                # grabbing "FY 2024" / "US 12"). e.g. "CAT320", "SRX340", "XPS 13",
                # "AX-760", "WS-C2960X".
                r'(?P<val>(?-i:\b(?:[A-Z]{2,4}-?\d{2,4}|[A-Z]{3,4}\s\d{2,4})(?:[\s-][A-Z0-9]{1,4})*\b))',
                # Legacy demo brand list (kept for template compatibility).
                r'(?P<val>\b(?:TechPro|DataMax|EcoLite|PowerGrid|CloudSuite|NetSys|SoftEdge|UltraGen|ProMax|EliteCore|LiteFlex|MaxPro|PlusCore|UltraMax|ProElite|LiteMax)\s*(?:Pro|Plus|Elite|Ultra|Lite|Max|Edition|Series|System|Suite|Platform|Solution)\w*\b)',
            ],
            "description": "Product names and model identifiers",
            "dummy_prefix": "PROD",
            "critical": True
        },
        "DIRECTOR_NAME": {
            "patterns": [
                r'(?P<label>(?<!\w)(?:director|manager|ceo|cto|cfo|president|head|name)\s*[:=#]?\s*)(?P<val>(?-i:[A-Z][a-z]+(?:[ ]+[A-Z][a-z]+){0,3}))',
                r'(?P<val>(?-i:(?:Mr\.|Ms\.|Mrs\.|Dr\.)\s+[A-Z][a-z]+[ ]+[A-Z][a-z]+))',
                r'(?P<val>(?-i:\b[A-Z][a-z]{1,4},[ ]+[A-Z][a-z]+(?:[ ]+[A-Z][a-z]+)*\b))',
                r'(?P<val>(?-i:\b(?:EG|AB|CD|EF|GH|IJ|KL|MN|OP|QR|ST|UV|WX|YZ)[A-Z]*[ ]+[A-Z][a-z]+\b))',
            ],
            "description": "Executive and personnel names",
            "dummy_prefix": "DIRECTOR",
            "critical": True
        },
        "QUOTATION_ID": {
            "patterns": [
                r'(?P<label>(?:quotation|quote|qtn|ref)\s*[:#]?\s*)(?P<val>[A-Z0-9-]+-\d{4,}[\w-]*)',
                r'(?P<val>\b[A-Z]{2,5}-\d{4}[\w-]*\b)',
            ],
            "description": "Quotation/reference identifiers",
            "dummy_prefix": "QUOTE",
            "critical": True
        },
        "COST_VALUE": {
            "patterns": [
                r'(?P<val>(?-i:\$[0-9]|€[0-9]|£[0-9]|RM\s?[0-9])\d{0,2}(?:[.,]\d{3})*(?:[.,]\d{2})?)',
                r'(?P<val>\b\d+(?:[.,]\d{2})?\s*(?-i:USD|EUR|MYR|GBP)\b)',
            ],
            "description": "Monetary values and costs",
            "dummy_prefix": "COST",
            "critical": True
        },
        "STANDARD_PART": {
            "patterns": [
                # Part NUMBER — only when explicitly labelled (Part No / P/N / Item No /
                # Material / Cat No / Model No / SKU / Stock Code). Never matches bare
                # tokens, so a lone catalogue ref is left alone unless it co-occurs
                # (on the same row) with a cost/price or a description -> clustered.
                r'(?P<label>(?:part\s*(?:no|number|num|#)?|p/?n|item(?:\s*no\.?)?|material|mat(?:erial)?\s*code|cat(?:alog)?\s*no\.?|model\s*no\.?|sku|stock\s*code)\s*[:#=]?\s*)(?P<val>(?-i:[A-Z0-9][A-Z0-9\-/.]{2,22}))',
                # Part DESCRIPTION — labelled (Description / Desc / Spec ONLY — bare
                # "part"/"item" are excluded to avoid colliding with Part No / Item No).
                # The value stops before a price keyword or another field label on the
                # same row, so it does not swallow the cost/price that follows it.
                r'(?P<label>(?:description|desc|spec(?:ification)?)\s*[:#=]?\s*)(?P<val>[A-Za-z0-9][\w\-/ ]*?)(?=\s+(?:unit\s*(?:price|cost)|price|cost|amount|total|part\s*no|p/?n|item\s*no|model\s*no|sku|rm|\$|€|£)\b|[,;|\n\r]|$)',
            ],
            "description": "Standard off-the-shelf (OTS/COTS) part info: number, description, unit cost",
            "dummy_prefix": "STDPART",
            "critical": True
        }
    },
    "HEALTH_INFORMATION": {
        "MEDICAL_RECORD_NUMBER": {
            "patterns": [
                r'(?P<label>(?:mrn|medical\s*record|patient\s*id)\s*[:#]?\s*)(?P<val>[A-Z0-9]{6,12})',
                r'(?P<val>\b\d{6,12}\b)(?=\D*(?:patient|medical|health))',
            ],
            "description": "Medical record numbers",
            "dummy_prefix": "MRN",
            "critical": True
        },
        "HEALTH_PLAN_BENEFICIARY": {
            "patterns": [
                r'(?P<label>(?:health\s*plan|insurance|beneficiary)\s*[:#]?\s*)(?P<val>[A-Z0-9]{6,12})',
            ],
            "description": "Health plan beneficiary numbers",
            "dummy_prefix": "HPBN",
            "critical": True
        },
        "CONDITION": {
            "patterns": [
                r'(?P<label>(?:diagnosed\s*with|suffering\s*from|history\s*of|diagnosis|condition)\s*[:#]?\s*)(?P<val>(?-i:[A-Z][a-z]+(?:[ ]+[A-Z0-9]?[a-z0-9]+){0,3}))',
                r'(?P<val>(?-i:\b(?:diabetes|hypertension|cancer|asthma|depression|anxiety|covid|hiv|leukemia)\b))',
            ],
            "description": "Medical conditions",
            "dummy_prefix": "COND",
            "critical": False
        }
    },
    "GOVERNMENT_IDS": {
        "PASSPORT": {
            "patterns": [
                r'(?P<label>(?:passport|passport\s*no\.?)\s*[:#]?\s*)(?P<val>(?-i:[A-Z0-9]{6,9}))',
                r'(?P<val>(?-i:\b[A-Z]{2}\d{6,7}\b))',
            ],
            "description": "Passport numbers",
            "dummy_prefix": "PASS",
            "critical": True
        },
        "DRIVER_LICENSE": {
            "patterns": [
                r'(?P<label>(?:driver\'?s?\s*license|dl)\s*[:#]?\s*)(?P<val>[A-Z0-9]{4,12})',
            ],
            "description": "Driver license numbers",
            "dummy_prefix": "DL",
            "critical": True
        },
        "TAX_ID": {
            "patterns": [
                r'(?P<label>(?:tax\s*id|tin|ein)\s*[:#]?\s*)(?P<val>\d{2}-?\d{7})',
            ],
            "description": "Tax identification numbers",
            "dummy_prefix": "TAX",
            "critical": True
        }
    },
    "FINANCIAL": {
        "IBAN": {
            "patterns": [
                r'(?P<label>(?:iban)\s*[:#]?\s*)(?P<val>(?-i:[A-Z]{2}\d{2}(?:[ ]?[A-Z0-9]){4,30}))',
                r'(?P<val>(?-i:\b[A-Z]{2}\d{2}(?:[ ]?[A-Z0-9]){4,30}\b))',
            ],
            "description": "International Bank Account Number",
            "dummy_prefix": "IBAN",
            "critical": True
        },
        "SWIFT": {
            "patterns": [
                r'(?P<label>(?:swift|bic)\s*[:#]?\s*)(?P<val>(?-i:[A-Z]{6}[A-Z0-9]{2,5}))',
                r'(?P<val>(?-i:\b[A-Z]{6}[A-Z0-9]{2,5}\b))',
            ],
            "description": "SWIFT/BIC codes",
            "dummy_prefix": "SWIFT",
            "critical": True
        },
        "BITCOIN_ADDRESS": {
            "patterns": [
                r'(?P<val>\b[13][a-km-zA-HJ-NP-Z1-9]{25,34}\b)',
                r'(?P<val>\bbc1[ac-hj-np-z02-9]{39,59}\b)',
            ],
            "description": "Bitcoin addresses",
            "dummy_prefix": "BTC",
            "critical": True
        }
    },
    "LOCATION_DATA": {
        "ADDRESS": {
            "patterns": [
                r'(?P<label>(?:address|addr)\s*[:#]?\s*)(?P<val>\d{1,5}\s+[A-Za-z0-9]+(?:\s+[A-Za-z0-9]+)*(?:\s+(?:Street|St\.|Avenue|Ave\.|Road|Rd\.|Boulevard|Blvd\.|Lane|Ln\.|Drive|Dr\.|Court|Ct\.|Place|Pl\.|Jalan|Bukit|Lorong|Taman|Persiaran)))',
                r'(?P<val>\d{1,5}\s+[A-Za-z0-9]+(?:\s+[A-Za-z0-9]+)*(?:\s+(?:Street|St\.|Avenue|Ave\.|Road|Rd\.|Boulevard|Blvd\.|Lane|Ln\.|Drive|Dr\.|Court|Ct\.|Place|Pl\.|Jalan|Bukit|Lorong|Taman|Persiaran)))',
            ],
            "description": "Street addresses",
            "dummy_prefix": "ADDR",
            "critical": False
        },
        "GPS_COORDINATES": {
            "patterns": [ r'(?P<val>-?\d{1,2}\.\d{4,6},\s*-?\d{1,3}\.\d{4,6})' ],
            "description": "GPS coordinates",
            "dummy_prefix": "GPS",
            "critical": False
        }
    },
    "CREDENTIALS": {
        "USERNAME": {
            "patterns": [
                r'(?P<label>(?:username|user\s*name|login|uid|user)\s*[:=#]?\s*)(?P<val>[a-zA-Z0-9_.-]{3,20})',
            ],
            "description": "Usernames",
            "dummy_prefix": "USER",
            "critical": True
        },
        "PASSWORD": {
            "patterns": [
                r'(?P<label>(?:password|passwd|pwd)\s*[:=#]?\s*)(?P<val>[^\s]{6,64})',
            ],
            "description": "Passwords",
            "dummy_prefix": "PASS",
            "critical": True
        },
        "API_KEY": {
            "patterns": [
                r'(?P<label>(?:api[_-]?key|apikey|secret[_-]?key)\s*[:=#]?\s*)(?P<val>[A-Za-z0-9_\-]{20,64})',
            ],
            "description": "API keys",
            "dummy_prefix": "KEY",
            "critical": True
        }
    }
}


# ─── Smart Dummy Generator ───────────────────────────────────────────────────
# Goal (user requirement): keep the DOCUMENT STRUCTURE intact (labels like
# "account number:", "Name:", "contact:" stay), and replace ONLY the sensitive
# VALUE with a realistic, format-preserving, structurally-valid dummy. The real
# value -> dummy value mapping is stored (reversible). Deterministic: the same
# real value always maps to the same dummy within a process run (seeded hash),
# so duplicates stay coherent and the mapping is trivially invertible.

_FIRST_NAMES = ["Abraham","Sarah","David","Maria","Wei","Fatimah","John","Priya","Chen","Olivia","Omar","Lena","Raj","Nadia","Lucas","Hana"]
_LAST_NAMES  = ["Tan","Lim","Kumar","Singh","Lee","Rahman","Wong","Fernandez","Ng","Andersson","Hassan","Novak","Cruz","Mensah","Petrov","Yamamoto"]
_FAKE_BANKS = ["RHB Bank","Maybank","CIMB Bank","Public Bank","AmBank","Hong Leong Bank","UOB Bank","OCBC Bank"]

def _seed(real_value):
    """Stable 32-bit seed from the real value (same value -> same seed)."""
    return int(hashlib.sha256(real_value.encode("utf-8")).hexdigest(), 16)

def _digits(n, rng):
    return "".join(str(rng.randint(0, 9)) for _ in range(n))

def _mask_middle(digits, rng, keep_head=2, keep_tail=2):
    """Preserve a few leading/trailing digits, dummy the middle, keep length."""
    if len(digits) <= keep_head + keep_tail:
        return _digits(len(digits), rng)
    head = digits[:keep_head]
    tail = digits[-keep_tail:]
    mid = "".join(str(rng.randint(0, 9)) for _ in range(len(digits) - keep_head - keep_tail))
    return head + mid + tail

# Per-category deterministic dummy generators. Each takes (original_text,
# matched_value, rng) and returns a safe, shape-preserving substitute.
def gen_account(orig, val, rng):
    digits = re.sub(r"\D", "", val)
    if not digits:
        return "X0000000000"
    # Keep same length; prefix a soft 'X' only if shorter than 9 to signal dummy
    fake = "X" + _mask_middle(digits[1:] if len(digits) > 1 else digits, rng, 1, 2) if len(digits) <= 11 else _mask_middle(digits, rng, 2, 3)
    return fake[:len(digits)+1] if digits[:1] != "X" else fake

def gen_ssn(orig, val, rng):
    nums = re.sub(r"\D", "", val)
    if len(nums) >= 9:
        return f"{_digits(3,rng)}-{_digits(2,rng)}-{_digits(4,rng)}"
    return "XXX-XX-" + _digits(4, rng)

def gen_email(orig, val, rng):
    local = re.split(r"[@]", val)[0]
    seed = _seed(val)
    r = random.Random(seed)
    name = _FIRST_NAMES[r.randint(0, len(_FIRST_NAMES)-1)].lower()
    return f"{name}.dummy{r.randint(100,999)}@example.com"

def gen_phone(orig, val, rng):
    digits = re.sub(r"\D", "", val)
    if len(digits) >= 7:
        # Keep last 4 of the real number (format-preserving, still fake-looking)
        return f"+700****{digits[-4:]}"
    return "+700****0000"

def gen_credit_card(orig, val, rng):
    digits = re.sub(r"\D", "", val)
    if len(digits) >= 15:
        return " ".join(_digits(4, rng) for _ in range((len(digits)+3)//4))[:len(val)]
    return _mask_middle(digits, rng, 2, 4)

def gen_name(orig, val, rng):
    # Keep the same structural capitalization/spacing, swap to a fake name.
    seed = _seed(val.strip())
    r = random.Random(seed)
    first = _FIRST_NAMES[r.randint(0, len(_FIRST_NAMES)-1)]
    last = _LAST_NAMES[r.randint(0, len(_LAST_NAMES)-1)]
    # Preserve "Last, First" vs "First Last" vs "TITLE First Last" shapes
    if "," in val:
        return f"{last}, {first}"
    if re.match(r"^(Mr|Ms|Mrs|Dr|Director|CEO|CTO|CFO|President|Head)\b", val, re.I):
        title = re.match(r"^(Mr|Ms|Mrs|Dr|Director|CEO|CTO|CFO|President|Head)\b\.?\s*", val, re.I).group(0)
        return f"{title}{first} {last}"
    return f"{first} {last}"

def gen_company(orig, val, rng):
    seed = _seed(val)
    r = random.Random(seed)
    suffixes = ["Corporation","Inc.","LLC","Ltd.","Group","Holdings","Sdn Bhd"]
    roots = ["Apex","Vertex","Nova","Summit","Helios","Quanta","Meridian","Atlas","Orion","Pinnacle"]
    return f"{roots[r.randint(0, len(roots)-1)]} {suffixes[r.randint(0, len(suffixes)-1)]}"

def gen_quote(orig, val, rng):
    m = re.match(r"^([A-Z]{2,5})-(\d{4})([-\w]*)$", val)
    if m:
        return f"{m.group(1)}-{_digits(4,rng)}{m.group(3)}"
    return "REF-" + _digits(6, rng)

def gen_cost(orig, val, rng):
    sym = ""
    if val[:1] in "$€£":
        sym = val[:1]
    elif val[:2].upper() == "RM":
        sym = "RM"
    nums = re.sub(r"[^0-9.,]", "", val)
    if "." in nums:
        intp, dec = nums.split(".", 1)
        fake = f"{_digits(max(1,len(intp)),rng)}.{dec[:2]}"
    else:
        fake = _digits(max(2, len(nums)), rng)
    return sym + fake

def gen_mrn(orig, val, rng):
    digits = re.sub(r"\D", "", val)
    if len(digits) >= 4:
        return "MRN " + _digits(len(digits), rng)
    return "MRN " + _digits(6, rng)

def gen_iban(orig, val, rng):
    m = re.match(r"^([A-Z]{2}\d{2})", val)
    prefix = m.group(1) if m else "XX00"
    body = re.sub(r"[^A-Z0-9]", "", val)[len(prefix):]
    fake_body = "".join(_digits(1, rng) if c.isdigit() else chr(ord('A')+rng.randint(0,25)) for c in body)
    # Re-insert the original spacing layout so the dummy keeps the same shape
    spaced = ""
    bi = 0
    for ch in val[len(prefix):]:
        if ch in " -":
            spaced += ch
        else:
            spaced += fake_body[bi]; bi += 1
    return prefix + spaced

def gen_swift(orig, val, rng):
    return "".join(chr(ord('A')+rng.randint(0,25)) for _ in range(4)) + _digits(2, rng) + "".join(chr(ord('A')+rng.randint(0,25)) for _ in range(2))

def gen_btc(orig, val, rng):
    return "1" + "".join(chr(ord('A')+rng.randint(0,25)) for _ in range(min(len(val)-1, 33)))

def gen_address(orig, val, rng):
    # Keep the street type (Street/Road/Jalan/...) exactly; fake only the
    # house number and the road name. Preserve trailing postcode/city tokens.
    m = re.match(r"^(\d{1,5})\s+(.+?)\s+(Street|St\.|Avenue|Ave\.|Road|Rd\.|Boulevard|Blvd\.|Lane|Ln\.|Drive|Dr\.|Court|Ct\.|Place|Pl\.|Jalan|Bukit|Lorong|Taman|Persiaran)\b(.*)$", val, re.I)
    if m:
        st = m.group(3)
        tail = m.group(4)  # e.g. ", Taman Damai, 56000 Kuala Lumpur"
        return f"{rng.randint(1,999)} {_LAST_NAMES[rng.randint(0,len(_LAST_NAMES)-1)]} {st}{tail}"
    # Fallback: keep any trailing postcode/town, fake the leading number+name
    m2 = re.match(r"^(\d{1,5})\s+(.+?)(\s+\d{5}.*)$", val)
    if m2:
        return f"{rng.randint(1,999)} {_LAST_NAMES[rng.randint(0,len(_LAST_NAMES)-1)]}{m2.group(3)}"
    return f"{rng.randint(1,999)} {_LAST_NAMES[rng.randint(0,len(_LAST_NAMES)-1)]} Road"

def gen_gps(orig, val, rng):
    # Keep structure, jitter coordinates
    parts = re.findall(r"-?\d+\.\d+", val)
    if len(parts) >= 2:
        lat = float(parts[0]); lon = float(parts[1])
        return f"{lat + rng.uniform(-0.01,0.01):.5f}, {lon + rng.uniform(-0.01,0.01):.5f}"
    return val

def gen_username(orig, val, rng):
    seed = _seed(val)
    r = random.Random(seed)
    return _FIRST_NAMES[r.randint(0, len(_FIRST_NAMES)-1)].lower() + str(r.randint(10,99))

def gen_password(orig, val, rng):
    return "********"  # passwords are never echoed as fake plaintext

def gen_apikey(orig, val, rng):
    return "AKIA" + _digits(16, rng)

def gen_condition(orig, val, rng):
    # Keep the same word-shape; swap to a dummy medical term of similar form.
    seed = _seed(val.strip())
    r = random.Random(seed)
    fake_terms = ["Hypertension","Asthma","Eczema","Migraine","Anemia","Bronchitis",
                  "Gastritis","Arthritis","Hypothyroid","Insomnia"]
    words = val.strip().split()
    out = []
    for w in words:
        if re.match(r"^\d+$", w):
            out.append(_digits(len(w), r))
        else:
            out.append(fake_terms[r.randint(0, len(fake_terms)-1)])
    return " ".join(out)

def gen_product(orig, val, rng):
    # Keep the SAME structural shape (Brand Model / SKU), swap to a fake brand+model.
    seed = _seed(val.strip())
    r = random.Random(seed)
    fake_brands = ["Acme","NovaTech","Vertex","Quantia","Helios","Orion","Pinnacle",
                   "Meridian","Atlas","Summit","Apex","Vantage","CoreLink","Stratus"]
    fake_models = ["Pro","Max","Ultra","Edge","Core","One","X","Series","Plus","Neo","Air","Mini"]
    # If it looks like a "<Brand> <model...>" (has a space), preserve that shape.
    parts = val.strip().split(None, 1)
    if len(parts) >= 2 and " " in val:
        brand = fake_brands[r.randint(0, len(fake_brands)-1)]
        model = parts[1]
        # Fake the model tail digits, keep letters/shape
        fake_model = re.sub(r"\d+", lambda m: _digits(len(m.group()), r), model)
        return f"{brand} {fake_model}"
    # Pure SKU like "XPS13" / "CAT320" -> fake brandless model
    body = re.sub(r"^[A-Z]{2,4}[-\s]?", "", val)
    fake_body = re.sub(r"\d+", lambda m: _digits(len(m.group()), r), body) if body else "100"
    return f"{fake_brands[r.randint(0,len(fake_brands)-1)]}{fake_body}"

def gen_generic(orig, val, rng):
    return "REDACTED_" + _digits(4, rng)

# Category -> generator mapping (fallback chain by category key / dummy_prefix)
DUMMY_GENERATORS = {
    "BANK_ACCOUNT": gen_account,
    "SSN": gen_ssn,
    "EMAIL": gen_email,
    "PHONE": gen_phone,
    "CREDIT_CARD": gen_credit_card,
    "DIRECTOR_NAME": gen_name,
    "COMPANY_NAME": gen_company,
    "PRODUCT_NAME": gen_product,
    "QUOTATION_ID": gen_quote,
    "COST_VALUE": gen_cost,
    "MEDICAL_RECORD_NUMBER": gen_mrn,
    "HEALTH_PLAN_BENEFICIARY": gen_mrn,
    "PASSPORT": lambda o,v,r: "A" + _digits(7, r),
    "DRIVER_LICENSE": lambda o,v,r: "D" + _digits(7, r),
    "TAX_ID": lambda o,v,r: _digits(2, r) + "-" + _digits(7, r),
    "IBAN": gen_iban,
    "SWIFT": gen_swift,
    "BITCOIN_ADDRESS": gen_btc,
    "ADDRESS": gen_address,
    "GPS_COORDINATES": gen_gps,
    "CONDITION": gen_condition,
    "USERNAME": gen_username,
    "PASSWORD": gen_password,
    "API_KEY": gen_apikey,
    "STANDARD_PART": lambda o,v,r: _stdpart_dispatch(o,v,r),
    "STDPART": lambda o,v,r: _stdpart_dispatch(o,v,r),
}

def _stdpart_dispatch(orig, val, rng):
    # Choose the right dummy shape for a clustered product-information field:
    # price-like value -> cost; long phrase -> description; else part number.
    if re.match(r"(?i)^[A-Za-z]{0,3}\s?\$|\$|€|£|rm\s|usd|eur|myr|gbp|\d+[.,]\d{2}", val.strip()):
        return gen_part_cost(orig, val, rng)
    if len(val.strip()) > 18:
        return gen_part_desc(orig, val, rng)
    return gen_part_no(orig, val, rng)


def gen_part_no(orig, val, rng):
    # Keep the structural shape of a part number: letters + digits + separators.
    letters = "".join(ch for ch in val if ch.isalpha())
    prefix = letters[:2].upper() if letters else "PN"
    body = re.sub(r"[^0-9]", "", val)
    if not body:
        body = _digits(5, rng)
    return prefix + "-" + _digits(max(4, len(body)), rng)

def gen_part_desc(orig, val, rng):
    # Generic, neutral component description (format-preserving shape only).
    return "Standard component assembly"

def gen_part_cost(orig, val, rng):
    # Preserve currency symbol/code, fake the amount realistically.
    sym = ""
    if val[:1] in "$€£":
        sym = val[:1]
    elif val[:2].upper() == "RM":
        sym = "RM "
    else:
        m = re.match(r"\$?\s*(\d+(?:[.,]\d+)?)\s*(USD|EUR|MYR|GBP)", val, re.IGNORECASE)
        sym = (m.group(2).upper() + " ") if m else "RM "
    return sym + f"{rng.randint(5,999)}.{_digits(2, rng)}"

def make_dummy(category, dummy_prefix, original, value):
    """Deterministic, format-preserving dummy for a matched value."""
    rng = random.Random(_seed(original))
    gen = DUMMY_GENERATORS.get(category) or DUMMY_GENERATORS.get(dummy_prefix) or gen_generic
    try:
        return gen(original, value, rng)
    except Exception:
        return gen_generic(original, value, rng)


# ─── Enhanced Redaction Engine ───────────────────────────────────────────────
class EnhancedRedactionEngine:
    """Enterprise-grade data anonymization engine with business-sensitive detection."""
    
    def __init__(self, settings=None):
        # Build flat pattern list from policy
        self.categories = OrderedDict()
        self.settings = settings if settings else get_default_settings()
        self._build_patterns()
    
    def _build_patterns(self):
        """Build regex patterns from security policy, respecting user settings."""
        # Reset fully so disabled categories are REMOVED (not merely skipped on
        # re-add). Without this, a category built once at startup stays active
        # forever even after the user disables it.
        self.categories = {}
        for group_name, group_settings in self.settings.get("categories", {}).items():
            if not group_settings.get("enabled", True):
                continue
            categories = SECURITY_POLICY.get(group_name, {})
            for cat_key, cat_info in categories.items():
                # Check if this category is enabled in settings
                sub_settings = group_settings.get("subcategories", {}).get(cat_key, {})
                if not sub_settings.get("enabled", True):
                    continue
                patterns = cat_info.get('patterns', [cat_info.get('pattern', '')])
                for p in patterns:
                    if p:
                        if cat_key not in self.categories:
                            self.categories[cat_key] = []
                        self.categories[cat_key].append((p, cat_info))

        # Add custom categories
        for custom in self.settings.get("custom", []):
            cat_key = custom.get("name", "CUSTOM")
            if not custom.get("enabled", True):
                continue
            pattern_str = custom.get("pattern", "")
            if pattern_str:
                cat_info = {
                    "description": custom.get("description", cat_key),
                    "dummy_prefix": custom.get("dummy_prefix", cat_key),
                    "critical": custom.get("critical", False),
                    "category": custom.get("category", "CUSTOM")
                }
                if cat_key not in self.categories:
                    self.categories[cat_key] = []
                self.categories[cat_key].append((pattern_str, cat_info))
    
    def _cluster_standard_parts(self, all_matches, text):
        """'Smart' co-occurrence pass (your design): treat a set of part-information
        fields as ONE product-information record when >=2 of
        {part number, part description, cost/price} appear on the same row/closely
        clustered. A bare part number or description is only redacted when it co-occurs
        with a cost/price (or another part field) on the same line — this is what stops
        normal sentences and lone catalogue refs from being falsely flagged.

        Fields that are part of such a cluster are relabelled category='STANDARD_PART'.
        Standalone COST_VALUE matches are left as-is (the dedicated cost category).
        """
        PART_FIELD_CATS = {"STANDARD_PART"}          # the number/description fields
        COST_CATS = {"COST_VALUE", "STANDARD_PART"}  # price-bearing matches
        PRICE_NEAR = re.compile(r"(?i)\b(?:unit\s*price|price|cost|amount|total|rm|\$|€|£|usd|eur|myr|gbp)\b")

        def same_row(a, b):
            # Same line if there's no newline between them, OR within ~60 chars.
            lo, hi = (a, b) if a['start'] < b['start'] else (b, a)
            between = text[lo['end']:hi['start']]
            if '\n' in between:
                return False
            return (hi['start'] - lo['end']) <= 60

        # Only look at STANDARD_PART field matches + COST_VALUE matches, plus any
        # PRODUCT_NAME SKU match that shares a row with a cost/part field (so a
        # labelled "Item: ABC-1234 ... RM 250.00" row is treated as one record).
        cand = [m for m in all_matches
                if m['category'] in PART_FIELD_CATS
                or m['category'] == "COST_VALUE"
                or m['category'] == "PRODUCT_NAME"]
        if not cand:
            return

        # Mark which candidates are "anchor" fields (part field or cost). A
        # PRODUCT_NAME match only joins a cluster if it is on the same row as an
        # anchor; standalone PRODUCT_NAME stays as PRODUCT_NAME.
        def is_anchor(m):
            return m['category'] in PART_FIELD_CATS or m['category'] == "COST_VALUE"

        # Find clusters: connected components where members are pairwise on the same row
        n = len(cand)
        parent = list(range(n))
        def find(x):
            while parent[x] != x:
                parent[x] = parent[parent[x]]; x = parent[x]
            return x
        def union(a, b):
            ra, rb = find(a), find(b)
            if ra != rb: parent[ra] = rb

        for i in range(n):
            for j in range(i + 1, n):
                if same_row(cand[i], cand[j]):
                    # Only connect a PRODUCT_NAME to the cluster if at least one of
                    # them is an anchor (part field or cost).
                    if cand[i]['category'] == "PRODUCT_NAME" and cand[j]['category'] == "PRODUCT_NAME":
                        continue
                    if is_anchor(cand[i]) or is_anchor(cand[j]):
                        union(i, j)

        # For each cluster, decide: it's a product-information record if it has
        # >=1 part field AND (>=1 cost/price OR >=2 part fields).
        groups = {}
        for i in range(n):
            groups.setdefault(find(i), []).append(cand[i])
        for root, members in groups.items():
            part_fields = [m for m in members if m['category'] == "STANDARD_PART"]
            costs = [m for m in members if m['category'] == "COST_VALUE" or PRICE_NEAR.search(m['text'])]
            is_record = len(part_fields) >= 1 and (len(costs) >= 1 or len(part_fields) >= 2)
            if is_record:
                for m in members:
                    # Relabel everything in the record as the unified product-info
                    # category (STANDARD_PART / STDPART), including a PRODUCT_NAME SKU
                    # that shares the row — so a labelled "Item: ABC-1234 ... RM 250"
                    # row is redacted as one coherent record.
                    m['category'] = "STANDARD_PART"
                    m['cat_info'] = dict(m['cat_info'])
                    if m['cat_info'].get('dummy_prefix') == "PROD":
                        m['cat_info']['description'] = "Standard part number"
                    else:
                        m['cat_info']['description'] = "Standard part (product information record)"
                    m['cat_info']['dummy_prefix'] = "STDPART"

    def redact(self, text, style=None):
        """Perform redaction with reversible mapping.

        style="token" (default): replace matches with opaque {PREFIX_n} tokens.
        style="smart": replace ONLY the sensitive VALUE with a realistic,
            format-preserving dummy (keeps document labels/structure). The
            real -> dummy mapping is recorded (reversible).
        """
        if style is None:
            style = self.settings.get("redaction_style", "token") if self.settings else "token"
        smart = (style == "smart")

        redactions = []
        category_counts = {}
        # For smart mode: cache so the same real value maps to the same dummy
        dummy_cache = {}

        # Build priority order dynamically from the SECURITY_POLICY group order
        built = set(self.categories.keys())
        priority_order = []
        for _grp, _cats in SECURITY_POLICY.items():
            for _cat in _cats:
                if _cat in built and _cat not in priority_order:
                    priority_order.append(_cat)
        for _cat in built:
            if _cat not in priority_order:
                priority_order.append(_cat)

        # Collect all matches across all categories
        all_matches = []
        for category in priority_order:
            if category not in self.categories:
                continue
            for pattern, cat_info in self.categories[category]:
                for match in re.finditer(pattern, text, re.IGNORECASE):
                    # The sensitive VALUE lives in the named capture group 'val'
                    # when the pattern uses one. An optional LABEL prefix (group
                    # 'label') is preserved as context. Patterns may instead use a
                    # single unnamed capture group for the value, or none at all.
                    gd = match.groupdict()
                    if gd.get('val') is not None:
                        value = gd['val']
                        vstart, vend = match.start('val'), match.end('val')
                    elif match.lastindex and match.lastindex >= 1:
                        value = match.group(1)
                        vstart, vend = match.start(1), match.end(1)
                    else:
                        value = match.group()
                        vstart, vend = match.start(), match.end()
                    all_matches.append({
                        'start': match.start(),
                        'end': match.end(),
                        'val_start': vstart,
                        'val_end': vend,
                        'text': match.group(),
                        'value': value,
                        'category': category,
                        'cat_info': cat_info
                    })

        # ── "Smart" co-occurrence clustering for STANDARD_PART ──────────────
        # Part number / description fields are only redacted when they appear on
        # the same row as a cost/price (or another part field). This makes the
        # category recognise product-information records "married together".
        self._cluster_standard_parts(all_matches, text)

        # Sort by position, then by length (longer matches first for overlapping)
        all_matches.sort(key=lambda m: (m['start'], -m['end']))
        
        # Process matches, handling overlaps by position
        redaction_map = []
        last_end = 0
        result_parts = []
        counter = 0
        real_to_dummy = {}  # smart mode: real value -> dummy value
        
        for match in all_matches:
            if match['start'] < last_end:
                continue
            
            result_parts.append(text[last_end:match['start']])
            
            category = match['category']
            category_counts[category] = category_counts.get(category, 0) + 1
            counter = category_counts[category]
            
            dummy_prefix = match['cat_info'].get('dummy_prefix', category)
            original_value = match['value'].strip()
            # Preserve any label prefix ("account number: ") so structure survives
            label_prefix = text[match['start']:match['val_start']]
            
            if smart:
                # Reuse dummy if we already mapped this exact real value
                if original_value in dummy_cache:
                    replacement = dummy_cache[original_value]
                else:
                    replacement = make_dummy(category, dummy_prefix, original_value, match['value'])
                    dummy_cache[original_value] = replacement
                real_to_dummy[original_value] = replacement
                # Keep label context + insert format-preserving dummy value
                result_parts.append(label_prefix + replacement)
                var_name = replacement
            else:
                var_name = f"{{{dummy_prefix}_{counter}}}"
                # Token mode: in token mode, we keep it simple - replace whole match
                # but preserve label prefix so the document still reads structurally.
                result_parts.append(label_prefix + var_name)
            
            redaction_map.append({
                'variable': var_name,
                'original': original_value,
                'category': category,
                'group': match['cat_info'].get('description', category),
                'position': match['start']
            })
            
            last_end = match['end']
            redactions.append({
                'category': category,
                'variable': var_name,
                'original': original_value,
                'length': len(match['value'])
            })
        
        result_parts.append(text[last_end:])
        safe_text = ''.join(result_parts)
        
        return safe_text, redaction_map, category_counts, redactions, (real_to_dummy if smart else None)

# ─── File Extractors ────────────────────────────────────────────────────────
def extract_text(file_path):
    """Extract text from various document formats using stdlib + optional libs."""
    path = Path(file_path)
    suffix = path.suffix.lower()
    
    try:
        if suffix in ('.txt', '.md', '.log', '.json', '.csv'):
            return path.read_text(encoding='utf-8', errors='replace')
        
        if suffix == '.pdf':
            try:
                sys.path.insert(0, VENV_SITE_PACKAGES)
                from pypdf import PdfReader
                reader = PdfReader(str(path))
                return '\n'.join([page.extract_text() or '' for page in reader.pages])
            except ImportError:
                raw = path.read_bytes()
                text = raw.decode('utf-8', errors='replace')
                text = re.sub(r'stream.*?endstream', '', text, flags=re.DOTALL)
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\(.*?\)', lambda m: m.group()[1:-1], text)
                return text
        
        if suffix == '.docx':
            try:
                sys.path.insert(0, VENV_SITE_PACKAGES)
                from docx import Document
                doc = Document(str(path))
                return '\n'.join([p.text for p in doc.paragraphs])
            except ImportError:
                try:
                    with zipfile.ZipFile(str(path)) as z:
                        for name in z.namelist():
                            if 'word/document' in name:
                                data_bytes = z.read(name)
                                text = re.sub(r'<[^>]+>', ' ', data_bytes.decode('utf-8', errors='ignore'))
                                text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text)
                                return text.strip()
                except Exception:
                    return f"DOCX extraction requires python-docx library"
        
        if suffix == '.xlsx':
            try:
                sys.path.insert(0, VENV_SITE_PACKAGES)
                from openpyxl import load_workbook
                wb = load_workbook(str(path), data_only=True)
                rows = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        rows.append(' '.join(str(c) if c else '' for c in row))
                return '\n'.join(rows)
            except ImportError:
                return f"XLSX extraction requires openpyxl library"
        
        if suffix == '.pptx':
            try:
                sys.path.insert(0, VENV_SITE_PACKAGES)
                from pptx import Presentation
                prs = Presentation(str(path))
                text_out = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text_out.append(shape.text)
                return '\n'.join(text_out)
            except ImportError:
                return f"PPTX extraction requires python-pptx library"
        
        if suffix in ('.html', '.htm'):
            return path.read_text(encoding='utf-8', errors='replace')
        
        # Fallback
        return path.read_text(encoding='utf-8', errors='replace')
    
    except Exception as e:
        return f"Error extracting text from {path.name}: {str(e)}"

# ─── HTML UI ─────────────────────────────────────────────────────────────────
HTML_UI = r"""
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>OAKAI Document Reader — Secure Data Anonymization</title>
<style>
  :root {
    --bg: #0b1220;
    --surface: #ffffff;
    --surface-2: #f8fafc;
    --surface-3: #eef2f7;
    --card-shadow: 0 12px 40px rgba(2, 8, 23, 0.28);
    --radius: 14px;
    --radius-lg: 20px;
    --primary: #0f172a;
    --accent: #0ea5e9;
    --accent-2: #6366f1;
    --text: #0f172a;
    --muted: #64748b;
    --border: #e2e8f0;
    --border-2: #cbd5e1;
    --ok: #10b981;
    --err: #ef4444;
    --warn: #f59e0b;
    --info: #3b82f6;
    --up-bg: #f1f5f9;
    --var-bg: #e0f2fe;
    --var-text: #0369a1;
    --grad-header: linear-gradient(120deg, #0f172a 0%, #1e293b 55%, #0ea5e9 140%);
  }
  * { margin: 0; padding: 0; box-sizing: border-box; }
  html, body { height: 100%; }
  body {
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, sans-serif;
    background:
      radial-gradient(1100px 500px at 12% -8%, rgba(14,165,233,0.18), transparent 60%),
      radial-gradient(900px 480px at 100% 0%, rgba(99,102,241,0.16), transparent 55%),
      var(--bg);
    min-height: 100vh;
    padding: 26px 18px 60px;
    color: var(--text);
  }
  .container { max-width: 1040px; margin: 0 auto; }

  .app-header {
    display: flex; align-items: center; justify-content: space-between;
    gap: 16px; flex-wrap: wrap;
    background: var(--grad-header);
    color: #fff; border-radius: var(--radius-lg);
    padding: 22px 26px; box-shadow: var(--card-shadow);
    position: relative; overflow: hidden;
  }
  .app-header::after {
    content: ""; position: absolute; inset: 0;
    background: radial-gradient(400px 200px at 90% 10%, rgba(255,255,255,0.08), transparent 70%);
    pointer-events: none;
  }
  .brand { display: flex; align-items: center; gap: 14px; }
  .brand-logo {
    width: 52px; height: 52px; border-radius: 14px;
    display: grid; place-items: center; font-size: 26px;
    background: rgba(255,255,255,0.14); border: 1px solid rgba(255,255,255,0.22);
    backdrop-filter: blur(6px);
  }
  .brand h1 { font-size: 25px; font-weight: 800; letter-spacing: 0.2px; }
  .brand h1 .accent { color: #7dd3fc; }
  .tagline { font-size: 13px; color: #cbd5e1; margin-top: 2px; }
  .header-actions { display: flex; gap: 10px; flex-wrap: wrap; }
  .btn {
    border: 1px solid transparent; border-radius: 10px; cursor: pointer;
    font-size: 14px; font-weight: 600; padding: 10px 16px; transition: .15s ease;
  }
  .btn-ghost {
    background: rgba(255,255,255,0.10); color: #fff; border-color: rgba(255,255,255,0.18);
  }
  .btn-ghost:hover { background: rgba(255,255,255,0.20); transform: translateY(-1px); }
  .btn-primary {
    background: linear-gradient(135deg, var(--accent), var(--accent-2));
    color: #fff; border: none; width: 100%; padding: 14px 22px;
    font-size: 16px; font-weight: 700; border-radius: 12px;
    box-shadow: 0 10px 24px rgba(14,165,233,0.35);
  }
  .btn-primary:hover:not(:disabled) { filter: brightness(1.06); transform: translateY(-1px); }
  .btn-primary:disabled { background: #94a3b8; box-shadow: none; cursor: not-allowed; }

  .stats-row { display: grid; grid-template-columns: repeat(4, 1fr); gap: 14px; margin: 18px 0; }
  .stat-card {
    background: var(--surface); border: 1px solid var(--border);
    border-radius: var(--radius); padding: 16px 18px; box-shadow: var(--card-shadow);
    position: relative; overflow: hidden;
  }
  .stat-card::before {
    content: ""; position: absolute; left: 0; top: 0; bottom: 0; width: 4px;
    background: linear-gradient(var(--accent), var(--accent-2));
  }
  .stat-card.secure::before { background: linear-gradient(var(--ok), #34d399); }
  .stat-value { font-size: 28px; font-weight: 800; color: var(--primary); line-height: 1; }
  .stat-label { font-size: 12px; color: var(--muted); margin-top: 6px; text-transform: uppercase; letter-spacing: .5px; }

  .card {
    background: var(--surface); border: 1px solid var(--border);
    border-radius: var(--radius-lg); padding: 22px; box-shadow: var(--card-shadow);
    margin-bottom: 18px;
  }
  .card-head { display: flex; align-items: center; justify-content: space-between; margin-bottom: 16px; }
  .card-head h2 { font-size: 18px; font-weight: 700; color: var(--primary); }
  .hint { font-size: 12px; color: var(--muted); }
  .link { color: var(--accent); font-weight: 600; }

  .upload-area {
    border: 2px dashed var(--border-2); border-radius: var(--radius);
    padding: 42px 24px; text-align: center; cursor: pointer;
    background: linear-gradient(180deg, #f8fafc, #f1f5f9); transition: .18s ease;
  }
  .upload-area:hover { border-color: var(--accent); background: #ecfeff; }
  .upload-area.drag-over { border-color: var(--accent-2); background: #eef2ff; transform: scale(1.005); }
  .upload-area input[type=file] { display: none; }
  .upload-icon { font-size: 44px; opacity: .8; }
  .upload-text { font-size: 16px; color: var(--text); margin-top: 8px; font-weight: 600; }
  .upload-sub { font-size: 12px; color: var(--muted); margin-top: 6px; }
  .file-info { font-size: 14px; color: var(--muted); margin: 12px 0; min-height: 20px; text-align: center; }

  .doc-grid { display: grid; grid-template-columns: 1fr; gap: 10px; }
  .doc-item {
    display: flex; align-items: center; justify-content: space-between; gap: 12px;
    padding: 14px 16px; background: var(--surface-2); border: 1px solid var(--border);
    border-radius: 12px; text-decoration: none; color: var(--text); transition: .15s;
  }
  .doc-item:hover { background: #eff6ff; border-color: var(--accent); transform: translateY(-1px); }
  .doc-filename { font-weight: 600; font-size: 14px; }
  .doc-meta { font-size: 12px; color: var(--muted); margin-top: 3px; }
  .doc-badge {
    font-size: 12px; font-weight: 700; color: #fff; background: linear-gradient(135deg, var(--accent), var(--accent-2));
    padding: 4px 12px; border-radius: 999px; white-space: nowrap;
  }
  .no-docs { text-align: center; color: var(--muted); padding: 28px; font-style: italic; }

  .result-header { display: flex; align-items: center; justify-content: space-between; gap: 12px; flex-wrap: wrap; }
  .result-title { font-size: 17px; font-weight: 700; }
  .status-badge { display: inline-flex; align-items: center; gap: 6px; padding: 6px 14px; border-radius: 999px; font-size: 13px; font-weight: 600; }
  .status-success { background: #dcfce7; color: #166534; }
  .status-error { background: #fee2e2; color: #991c1c; }
  .redaction-summary { display: flex; flex-wrap: wrap; gap: 8px; margin: 14px 0; }
  .redaction-badge {
    display: inline-flex; align-items: center; gap: 6px; padding: 5px 12px;
    border-radius: 999px; font-size: 12px; font-weight: 700; color: #fff;
    border-left: 4px solid rgba(255,255,255,0.55);
  }
  .info-table { width: 100%; border-collapse: collapse; font-size: 13px; margin: 6px 0 4px; }
  .info-table td { padding: 6px 8px; border-bottom: 1px solid var(--border); }
  .info-table td:first-child { color: var(--muted); width: 38%; }
  .info-table td:last-child { font-weight: 600; word-break: break-all; }
  .result-actions { display: flex; gap: 10px; flex-wrap: wrap; margin: 12px 0; }
  .act-btn {
    border: 1px solid var(--border-2); background: var(--surface-2); color: var(--primary);
    border-radius: 10px; padding: 9px 14px; font-size: 13px; font-weight: 600; cursor: pointer; transition: .15s;
    text-decoration: none; display: inline-block;
  }
  .act-btn:hover { background: var(--surface-3); border-color: var(--accent); }
  .result-content {
    background: #0f172a; color: #e2e8f0; border-radius: 12px; padding: 16px;
    font-family: 'SF Mono', 'Fira Code', Monaco, Consolas, monospace; font-size: 13px;
    line-height: 1.6; white-space: pre-wrap; word-break: break-word; max-height: 420px; overflow-y: auto; margin-top: 8px;
  }
  .variable-highlight { background: #0ea5e9; color: #04141f; padding: 1px 6px; border-radius: 5px; font-weight: 700; }

  .legend { display: flex; flex-wrap: wrap; gap: 8px; margin: 12px 0; }
  .legend-pill {
    display: inline-flex; align-items: center; gap: 7px; padding: 6px 12px; border-radius: 999px;
    font-size: 12px; font-weight: 700; cursor: pointer; border: 1px solid var(--border-2);
    background: var(--surface-2); color: var(--primary); transition: .15s; user-select: none;
  }
  .legend-pill .dot { width: 10px; height: 10px; border-radius: 50%; }
  .legend-pill.active { color: #fff; border-color: transparent; }
  .legend-pill:not(.active) { opacity: .7; }

  .mapping-section { margin-top: 16px; padding-top: 14px; border-top: 1px solid var(--border); }
  .mapping-section h4 { font-size: 13px; color: var(--muted); margin-bottom: 10px; }
  .mapping-table { width: 100%; border-collapse: collapse; font-size: 13px; }
  .mapping-table th, .mapping-table td { text-align: left; padding: 8px 10px; border-bottom: 1px solid var(--border); }
  .mapping-table th { color: var(--muted); font-weight: 600; }
  .mapping-table td:first-child { font-family: 'SF Mono', monospace; font-weight: 700; color: var(--accent); }
  .cat-tag { display: inline-flex; align-items: center; gap: 6px; padding: 3px 10px; border-radius: 999px; font-size: 11px; font-weight: 700; color: #fff; }
  .cat-tag .dot { width: 8px; height: 8px; border-radius: 50%; background: rgba(255,255,255,0.7); }
  .crit { font-size: 10px; padding: 2px 7px; border-radius: 5px; background: #fef3c7; color: #92400e; margin-left: 8px; font-weight: 700; }

  .g-PII { --gc:#ef4444; }
  .g-BUSINESS_SENSITIVE { --gc:#a855f7; }
  .g-HEALTH_INFORMATION { --gc:#14b8a6; }
  .g-GOVERNMENT_IDS { --gc:#3b82f6; }
  .g-FINANCIAL { --gc:#f59e0b; }
  .g-LOCATION_DATA { --gc:#6366f1; }
  .g-CREDENTIALS { --gc:#ec4899; }

  .modal-overlay {
    position: fixed; inset: 0; background: rgba(2,8,23,0.6); backdrop-filter: blur(3px);
    display: none; align-items: flex-start; justify-content: center; z-index: 1000; padding: 26px 16px; overflow-y: auto;
  }
  .modal-overlay.active { display: flex; }
  .modal {
    background: var(--surface); border-radius: var(--radius-lg); width: 100%; max-width: 860px;
    max-height: 92vh; overflow: hidden; box-shadow: 0 30px 80px rgba(2,8,23,0.5); display: flex; flex-direction: column;
  }
  .modal-head { padding: 20px 26px; border-bottom: 1px solid var(--border); display: flex; align-items: center; justify-content: space-between; }
  .modal-head h2 { font-size: 19px; font-weight: 700; color: var(--primary); }
  .modal-close { font-size: 24px; cursor: pointer; color: var(--muted); line-height: 1; border: none; background: none; }
  .modal-body { padding: 20px 26px; overflow-y: auto; }
  .modal-foot { padding: 16px 26px; border-top: 1px solid var(--border); display: flex; justify-content: flex-end; gap: 12px; }

  .settings-toolbar { display: flex; gap: 10px; margin-bottom: 16px; flex-wrap: wrap; }
  .search-input {
    flex: 1; min-width: 200px; padding: 11px 14px; border: 1px solid var(--border-2); border-radius: 10px;
    font-size: 14px; background: var(--surface-2);
  }
  .search-input:focus { outline: 2px solid var(--accent); border-color: var(--accent); }
  .mini-btn { padding: 9px 14px; border-radius: 10px; border: 1px solid var(--border-2); background: var(--surface-2); font-size: 13px; font-weight: 600; cursor: pointer; color: var(--primary); }
  .mini-btn:hover { border-color: var(--accent); background: #eff6ff; }

  .group {
    border: 1px solid var(--border); border-left: 5px solid var(--gc, #94a3b8);
    border-radius: 12px; margin-bottom: 14px; overflow: hidden; background: var(--surface-2);
  }
  .group-head {
    display: flex; align-items: center; justify-content: space-between; gap: 10px;
    padding: 14px 16px; cursor: pointer; background: var(--surface);
  }
  .group-title { display: flex; align-items: center; gap: 10px; font-weight: 700; font-size: 15px; color: var(--primary); }
  .group-dot { width: 12px; height: 12px; border-radius: 50%; background: var(--gc, #94a3b8); }
  .group-count { font-size: 12px; color: var(--muted); font-weight: 600; }
  .group-tools { display: flex; gap: 8px; align-items: center; }
  .chev { transition: transform .2s; color: var(--muted); font-size: 13px; }
  .group.collapsed .chev { transform: rotate(-90deg); }
  .group-body { padding: 6px 16px 14px; display: grid; gap: 8px; }
  .group.collapsed .group-body { display: none; }
  .cat-row {
    display: flex; align-items: center; justify-content: space-between; gap: 12px;
    padding: 11px 14px; background: #fff; border: 1px solid var(--border); border-radius: 10px;
  }
  .cat-row label { flex: 1; cursor: default; font-size: 14px; display: flex; align-items: center; gap: 8px; pointer-events: none; }
  .cat-row { cursor: pointer; }
  .cat-row input[type=checkbox] { width: 19px; height: 19px; accent-color: var(--gc, #0ea5e9); cursor: pointer; }
  .switch { position: relative; width: 42px; height: 24px; }
  .switch input { display: none; }
  .slider { position: absolute; inset: 0; background: #cbd5e1; border-radius: 999px; transition: .2s; cursor: pointer; }
  .slider::before { content: ""; position: absolute; width: 18px; height: 18px; left: 3px; top: 3px; background: blue; border-radius: 50%; transition: .2s; }
  .switch input:checked + .slider { background: var(--gc, #0ea5e9); }
  .switch input:checked + .slider::before { transform: translateX(18px); background: #fff; }

  .style-toggle-card { display: flex; align-items: center; justify-content: space-between; gap: 16px; background: linear-gradient(135deg, #f0f9ff, #e0f2fe); border: 1px solid #bae6fd; border-radius: 12px; padding: 14px 16px; margin: 4px 0 14px; }
  .style-toggle-title { font-size: 15px; font-weight: 800; color: #0369a1; margin-bottom: 4px; }
  .style-toggle-sub { font-size: 12px; color: #0c4a6e; line-height: 1.5; }
  .style-toggle-sub code { background: #fff; border: 1px solid #bae6fd; border-radius: 5px; padding: 1px 5px; font-size: 11px; color: #075985; }
  .switch.big { width: 56px; height: 32px; flex: 0 0 auto; }
  .switch.big .slider::before { width: 24px; height: 24px; left: 4px; top: 4px; background: #fff; }
  .switch.big input:checked + .slider { background: #0ea5e9; }
  .switch.big input:checked + .slider::before { transform: translateX(24px); }

  .custom-section { border-top: 2px dashed var(--border-2); margin-top: 14px; padding-top: 18px; }
  .custom-row { display: flex; gap: 8px; align-items: center; margin-bottom: 8px; flex-wrap: wrap; }
  .custom-row input, .custom-row select { flex: 1; min-width: 140px; padding: 9px 12px; border: 1px solid var(--border-2); border-radius: 8px; font-size: 13px; }
  .custom-row button { padding: 9px 16px; background: var(--accent); color: #fff; border: none; border-radius: 8px; font-size: 13px; font-weight: 600; cursor: pointer; }
  .custom-item { display: flex; align-items: center; justify-content: space-between; gap: 10px; padding: 9px 12px; background: #fff; border: 1px solid var(--border); border-radius: 8px; margin-bottom: 6px; }
  .custom-item button { background: var(--err); color: #fff; border: none; border-radius: 6px; padding: 6px 12px; cursor: pointer; font-size: 12px; }
  .banner-ok { background: var(--ok); color: #fff; padding: 12px 18px; border-radius: 10px; text-align: center; font-weight: 600; }
  .banner-err { background: #fef2f2; color: var(--err); padding: 10px 14px; border-radius: 8px; border: 1px solid #fee2e2; margin-bottom: 12px; font-size: 13px; }

  .app-footer { text-align: center; color: #94a3b8; font-size: 12px; margin-top: 24px; }

  @media (max-width: 720px) {
    .stats-row { grid-template-columns: repeat(2, 1fr); }
    .app-header { flex-direction: column; align-items: flex-start; }
  }
</style>
</head>
<body>
<div class="container">

  <header class="app-header">
    <div class="brand">
      <div class="brand-logo">🔒</div>
      <div>
        <h1>OAKAI <span class="accent">Document Reader</span></h1>
        <p class="tagline">Enterprise Data Anonymization &amp; PII Redaction Engine</p>
      </div>
    </div>
    <div class="header-actions">
      <button class="btn btn-ghost" onclick="openSettings()">⚙️ Settings</button>
      <button class="btn btn-ghost" onclick="loadDocs()">🔄 Refresh</button>
      <button class="btn btn-ghost" onclick="resetServer()">♻️ Reset Server</button>
    </div>
  </header>

  <div class="stats-row">
    <div class="stat-card"><div class="stat-value" id="statDocs">0</div><div class="stat-label">Documents</div></div>
    <div class="stat-card"><div class="stat-value" id="statRedactions">0</div><div class="stat-label">Redactions</div></div>
    <div class="stat-card"><div class="stat-value" id="statCats">0</div><div class="stat-label">Categories</div></div>
    <div class="stat-card secure"><div class="stat-value">100%</div><div class="stat-label">Local</div></div>
  </div>

  <section class="card">
    <div class="card-head"><h2>📥 Upload &amp; Process</h2><span class="hint">Drag &amp; drop or click to browse</span></div>
    <div class="upload-area" id="dropZone">
      <input type="file" id="fileInput" accept=".pdf,.docx,.xlsx,.pptx,.txt,.csv,.html,.htm,.md,.json" multiple>
      <div class="upload-icon">📎</div>
      <p class="upload-text">Drop files here or <span class="link">browse</span></p>
      <p class="upload-sub">PDF · DOCX · XLSX · PPTX · TXT · CSV · HTML · JSON · MD</p>
    </div>
    <div class="file-info" id="fileName">No file selected</div>
    <button class="btn-primary" id="processBtn" disabled>Upload &amp; Process</button>
  </section>

  <section class="card">
    <div class="card-head"><h2>📄 Processed Documents</h2><span class="doc-count" id="docCount"></span></div>
    <div id="docsList" class="doc-grid"><div class="no-docs">Loading documents…</div></div>
  </section>

  <section class="card" id="resultsSection" style="display:none;">
    <div class="result-header">
      <span class="result-title" id="resultTitle">Processing Results</span>
      <span class="status-badge status-success" id="resultStatus">✓ Complete</span>
    </div>
    <div class="redaction-summary" id="redactionSummary"></div>
    <table class="info-table">
      <tr><td>Original file</td><td id="resultFilename"></td></tr>
      <tr><td>Total redactions</td><td id="resultRedactions"></td></tr>
      <tr><td>Categories found</td><td id="resultCategories"></td></tr>
      <tr><td>Safe document</td><td id="resultSafePath"></td></tr>
      <tr><td>Variable mapping</td><td id="resultMapPath"></td></tr>
    </table>
    <div class="result-actions">
      <button class="act-btn" onclick="copySafe()">📋 Copy safe text</button>
      <a class="act-btn" id="downloadSafe" href="#" target="_blank">⬇️ Download safe JSON</a>
    </div>
    <div class="result-content" id="resultContent"></div>
    <div class="mapping-section">
      <h4>Redaction Variable Mapping (stored locally):</h4>
      <div class="legend" id="legend"></div>
      <table class="mapping-table" id="mappingTable">
        <thead><tr><th>Variable</th><th>Original Value</th><th>Category</th></tr></thead>
        <tbody id="mappingBody"></tbody>
      </table>
    </div>
  </section>

  <footer class="app-footer">OAKAI Document Reader v2.0 — All processing is 100% local. No APIs or tokens required.</footer>
</div>

<div class="modal-overlay" id="settingsOverlay">
  <div class="modal">
    <div class="modal-head">
      <h2>🔒 Redaction Settings</h2>
      <button class="modal-close" onclick="closeSettings()">×</button>
    </div>
    <div class="modal-body">
      <p style="color:var(--muted); margin-bottom:14px; font-size:14px;">
        Toggle which data categories to redact. Changes apply to the next document you process.
      </p>
      <div id="settingsError" class="banner-err" style="display:none;"></div>
      <div id="settingsSuccess" class="banner-ok" style="display:none;">✅ Settings saved successfully!</div>
      <div class="settings-toolbar">
        <input type="text" id="catSearch" class="search-input" placeholder="🔍 Search categories (e.g. passport, email, iban)…" oninput="filterCats()">
        <button class="mini-btn" onclick="setAll(true)">Enable all</button>
        <button class="mini-btn" onclick="setAll(false)">Disable all</button>
      </div>
      <div class="style-toggle-card">
        <div class="style-toggle-info">
          <div class="style-toggle-title">🪄 Smart Dummy Mode</div>
          <div class="style-toggle-sub">Replace only the <b>sensitive value</b> with a realistic, format-preserving dummy — keeps labels &amp; document structure intact (e.g. <code>account: 97652345334</code> → <code>account: X0003455334</code>). Mapping is stored and reversible. <b>Token mode</b> (opaque <code>{PREFIX_n}</code>) is the default.</div>
        </div>
        <label class="switch big"><input type="checkbox" id="smartToggle" onchange="setStyle(this.checked)"><span class="slider"></span></label>
      </div>
      <div id="settingsList"><div style="text-align:center; color:var(--muted); padding:20px;">Loading categories…</div></div>
      <div class="custom-section">
        <h3 style="font-size:14px; font-weight:700; margin-bottom:12px; color:var(--primary);">➕ Add Custom Category</h3>
        <div class="custom-row">
          <input type="text" id="customName" placeholder="Category name (e.g., PATENT_ID)">
          <input type="text" id="customPattern" placeholder="Regex pattern">
        </div>
        <div class="custom-row">
          <input type="text" id="customDesc" placeholder="Description">
          <select id="customPrefix">
            <option value="CUSTOM">CUSTOM</option>
            <option value="SSN">SSN</option>
            <option value="ID">ID</option>
            <option value="REF">REF</option>
            <option value="KEY">KEY</option>
          </select>
        </div>
        <button onclick="addCustomCategory()" style="width:100%; padding:10px; background:var(--accent); color:#fff; border:none; border-radius:8px; font-weight:600; cursor:pointer;">Add Custom Category</button>
      </div>
    </div>
    <div class="modal-foot">
      <button class="btn btn-ghost" style="color:var(--primary); border-color:var(--border-2);" onclick="closeSettings()">Cancel</button>
      <button class="btn-primary" style="width:auto;" onclick="saveSettings()">💾 Save Settings</button>
    </div>
  </div>
</div>

<script>
const dropZone = document.getElementById('dropZone');
const fileInput = document.getElementById('fileInput');
const fileNameDiv = document.getElementById('fileName');
const processBtn = document.getElementById('processBtn');
const resultsSection = document.getElementById('resultsSection');
const docCountSpan = document.getElementById('docCount');
let selectedFile = null;
let currentSettings = null;

const GROUP_META = {
  'PII': { label: 'PII / Sensitive', cls: 'g-PII' },
  'BUSINESS_SENSITIVE': { label: 'Business Sensitive', cls: 'g-BUSINESS_SENSITIVE' },
  'HEALTH_INFORMATION': { label: 'Health Information', cls: 'g-HEALTH_INFORMATION' },
  'GOVERNMENT_IDS': { label: 'Government IDs', cls: 'g-GOVERNMENT_IDS' },
  'FINANCIAL': { label: 'Financial Data', cls: 'g-FINANCIAL' },
  'LOCATION_DATA': { label: 'Location Data', cls: 'g-LOCATION_DATA' },
  'CREDENTIALS': { label: 'Credentials', cls: 'g-CREDENTIALS' }
};
const GROUP_COLOR = {
  'PII': '#ef4444', 'BUSINESS_SENSITIVE': '#a855f7', 'HEALTH_INFORMATION': '#14b8a6',
  'GOVERNMENT_IDS': '#3b82f6', 'FINANCIAL': '#f59e0b', 'LOCATION_DATA': '#6366f1', 'CREDENTIALS': '#ec4899'
};
const CATEGORY_GROUPS = {
  'SSN': 'PII', 'EMAIL': 'PII', 'PHONE': 'PII', 'CREDIT_CARD': 'PII', 'BANK_ACCOUNT': 'PII',
  'COMPANY_NAME': 'BUSINESS_SENSITIVE', 'PRODUCT_NAME': 'BUSINESS_SENSITIVE', 'DIRECTOR_NAME': 'BUSINESS_SENSITIVE', 'QUOTATION_ID': 'BUSINESS_SENSITIVE', 'COST_VALUE': 'BUSINESS_SENSITIVE',
  'MEDICAL_RECORD_NUMBER': 'HEALTH_INFORMATION', 'HEALTH_PLAN_BENEFICIARY': 'HEALTH_INFORMATION', 'CONDITION': 'HEALTH_INFORMATION',
  'PASSPORT': 'GOVERNMENT_IDS', 'DRIVER_LICENSE': 'GOVERNMENT_IDS', 'TAX_ID': 'GOVERNMENT_IDS',
  'IBAN': 'FINANCIAL', 'SWIFT': 'FINANCIAL', 'BITCOIN_ADDRESS': 'FINANCIAL',
  'ADDRESS': 'LOCATION_DATA', 'GPS_COORDINATES': 'LOCATION_DATA',
  'USERNAME': 'CREDENTIALS', 'PASSWORD': 'CREDENTIALS', 'API_KEY': 'CREDENTIALS'
};
function groupOf(cat){ return CATEGORY_GROUPS[cat] || 'PII'; }
function escapeHtml(s){ return (s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;'); }

dropZone.addEventListener('click', () => fileInput.click());
fileInput.addEventListener('change', () => {
  const files = fileInput.files;
  if (files.length > 0) {
    selectedFile = files[0];
    fileNameDiv.textContent = files.length === 1
      ? selectedFile.name + ' (' + Math.round(selectedFile.size/1024) + ' KB)'
      : files.length + ' files selected';
    processBtn.disabled = false;
    processBtn.textContent = 'Process ' + files.length + ' file(s)';
  } else { resetFileUI(); }
});
function resetFileUI(){
  fileNameDiv.textContent = 'No file selected';
  processBtn.disabled = true; processBtn.textContent = 'Upload & Process'; selectedFile = null;
}

processBtn.addEventListener('click', async () => {
  if (!selectedFile) return;
  processBtn.disabled = true; processBtn.textContent = 'Processing…';
  resultsSection.style.display = 'none';
  const resultContent = document.getElementById('resultContent');
  const resultTitle = document.getElementById('resultTitle');
  const resultStatus = document.getElementById('resultStatus');
  const formData = new FormData();
  formData.append('file', selectedFile);
  try {
    const res = await fetch('/upload', { method: 'POST', body: formData });
    const data = await res.json();
    processBtn.disabled = false; processBtn.textContent = 'Upload & Process';
    if (data.error) {
      resultTitle.textContent = 'Processing Error';
      resultStatus.className = 'status-badge status-error';
      resultStatus.textContent = '✗ Failed';
      resultContent.innerHTML = '<div class="banner-err">Error: ' + escapeHtml(data.error) + '</div>';
      resultsSection.style.display = 'block';
    } else {
      resultTitle.textContent = data.original_filename || 'Document Processed';
      resultStatus.className = 'status-badge status-success';
      resultStatus.textContent = '✓ Complete';
      document.getElementById('resultFilename').textContent = data.original_filename || 'N/A';
      document.getElementById('resultRedactions').textContent = data.total_redactions;
      document.getElementById('resultCategories').textContent = Object.keys(data.category_counts || {}).join(', ');
      document.getElementById('resultSafePath').textContent = data.safe_url || 'N/A';
      document.getElementById('resultMapPath').textContent = 'stored at: ' + (data.map_path || 'data/redaction_maps/');
      document.getElementById('downloadSafe').href = data.safe_url || '#';

      let badges = '';
      for (const [cat, count] of Object.entries(data.category_counts || {})) {
        const g = groupOf(cat); const col = GROUP_COLOR[g] || '#64748b';
        badges += '<span class="redaction-badge" style="background:' + col + '">' + cat + ': ' + count + '</span>';
      }
      document.getElementById('redactionSummary').innerHTML = badges || '<span class="hint">No sensitive data detected.</span>';

      try {
        const safeRes = await fetch('/documents/' + data.document_id + '/safe');
        const safeData = await safeRes.json();
        let safeText = safeData.all_text || 'No text extracted';
        safeText = escapeHtml(safeText);
        safeText = safeText.replace(/\{([A-Z]+_\d+)\}/g, '<span class="variable-highlight">{$1}</span>');
        resultContent.innerHTML = safeText;
      } catch(e) { resultContent.innerHTML = '<div class="banner-err">Could not load safe document content</div>'; }

      try {
        const mapRes = await fetch('/documents/' + data.document_id + '/map');
        const mapData = await mapRes.json();
        const mapping = mapData.full_map || {};
        let groupsPresent = new Set();
        let rows = '';
        for (const [varName, info] of Object.entries(mapping)) {
          const g = groupOf(info.category); const col = GROUP_COLOR[g] || '#64748b';
          groupsPresent.add(g);
          const crit = info.critical ? '<span class="crit">CRITICAL</span>' : '';
          rows += '<tr data-group="' + g + '">'
            + '<td>' + escapeHtml(varName) + '</td>'
            + '<td style="font-family:monospace;">' + escapeHtml(info.original) + '</td>'
            + '<td><span class="cat-tag" style="background:' + col + '"><span class="dot"></span>' + info.category + '</span>' + crit
            + '<br><small style="color:var(--muted);">' + escapeHtml(info.group) + '</small></td></tr>';
        }
        document.getElementById('mappingBody').innerHTML = rows || '<tr><td colspan="3">No variables mapped</td></tr>';
        buildLegend(groupsPresent);
      } catch(e2) { document.getElementById('mappingBody').innerHTML = '<tr><td colspan="3">Mapping unavailable</td></tr>'; }
      resultsSection.style.display = 'block';
    }
    resetFileUI();
    fileInput.value = '';
    await loadDocs();
  } catch(e) {
    processBtn.disabled = false; processBtn.textContent = 'Upload & Process';
    resultTitle.textContent = 'Error';
    resultStatus.className = 'status-badge status-error';
    resultStatus.textContent = '✗ Failed';
    resultContent.innerHTML = '<div class="banner-err">Error: ' + escapeHtml(e.message) + '</div>';
    resultsSection.style.display = 'block';
  }
});

function buildLegend(groupsPresent){
  const legend = document.getElementById('legend');
  legend.innerHTML = '';
  for (const g of groupsPresent) {
    const col = GROUP_COLOR[g] || '#64748b';
    const pill = document.createElement('span');
    pill.className = 'legend-pill active';
    pill.style.background = col;
    pill.innerHTML = '<span class="dot" style="background:rgba(255,255,255,0.8);"></span>' + (GROUP_META[g] ? GROUP_META[g].label : g);
    pill.onclick = () => { pill.classList.toggle('active'); filterMapping(); };
    legend.appendChild(pill);
  }
}
function filterMapping(){
  const pills = document.querySelectorAll('#legend .legend-pill');
  const active = new Set();
  pills.forEach(p => { if (p.classList.contains('active')) active.add(p.textContent.trim()); });
  document.querySelectorAll('#mappingBody tr[data-group]').forEach(tr => {
    tr.style.display = active.has(tr.getAttribute('data-group')) ? '' : 'none';
  });
}

async function copySafe(){
  const txt = document.getElementById('resultContent').innerText;
  try { await navigator.clipboard.writeText(txt); alert('Safe text copied to clipboard.'); }
  catch(e){ alert('Copy failed: ' + e.message); }
}

async function resetServer(){
  if (!confirm('Restart the server now? This picks up any changes made to doc_reader_onefile.py. The page reloads automatically.')) return;
  try { await fetch('/restart', { method: 'POST' }); } catch(e) {}
  const s = document.createElement('div');
  s.textContent = 'Restarting server, please wait…';
  s.style.cssText = 'position:fixed;top:20px;left:50%;transform:translateX(-50%);background:#0f172a;color:#fff;padding:12px 24px;border-radius:8px;z-index:9999;font-size:14px;';
  document.body.appendChild(s);
  let attempts = 0;
  const tryReload = async () => {
    attempts++;
    try { const r = await fetch('/health', { cache: 'no-store' }); if (r.ok) { location.reload(); return; } }
    catch(e) {}
    if (attempts < 20) setTimeout(tryReload, 500);
    else s.textContent = 'Still restarting… refresh manually.';
  };
  setTimeout(tryReload, 1500);
}

async function loadDocs(){
  try {
    const res = await fetch('/documents');
    const data = await res.json();
    const docs = data.documents || [];
    document.getElementById('statDocs').textContent = docs.length;
    let totalRed = 0;
    docs.forEach(d => { totalRed += (d.redactions_count||0); });
    docCountSpan.textContent = '(' + docs.length + ')';
    const list = document.getElementById('docsList');
    if (!docs.length) { list.innerHTML = '<div class="no-docs">No documents processed yet. Upload your first document above!</div>'; return; }
    let html = '';
    docs.slice(0,12).forEach(doc => {
      const d = doc.id.replace('doc_','').replace(/_/g,' ');
      html += '<a href="' + doc.url + '" target="_blank" class="doc-item">'
        + '<div><div class="doc-filename">' + escapeHtml(doc.filename) + '</div>'
        + '<div class="doc-meta">' + d + ' · ' + doc.size + ' bytes</div></div>'
        + '<span class="doc-badge">' + (doc.redactions_count||'?') + ' redactions</span></a>';
    });
    list.innerHTML = html;
    document.getElementById('statRedactions').textContent = totalRed;
    document.getElementById('statCats').textContent = Object.keys(GROUP_META).length;
  } catch(e) { document.getElementById('docsList').innerHTML = '<div class="no-docs">Error loading documents</div>'; }
}

['dragover','dragenter'].forEach(ev => dropZone.addEventListener(ev, e => { e.preventDefault(); dropZone.classList.add('drag-over'); }));
['dragleave','drop'].forEach(ev => dropZone.addEventListener(ev, e => { e.preventDefault(); dropZone.classList.remove('drag-over'); }));
dropZone.addEventListener('drop', e => {
  if (e.dataTransfer.files.length) {
    fileInput.files = e.dataTransfer.files;
    selectedFile = e.dataTransfer.files[0];
    fileNameDiv.textContent = selectedFile.name + ' (' + Math.round(selectedFile.size/1024) + ' KB)';
    processBtn.disabled = false; processBtn.textContent = 'Process ' + e.dataTransfer.files.length + ' file(s)';
    if (e.dataTransfer.files.length === 1) setTimeout(() => processBtn.click(), 100);
  }
});

async function openSettings(){
  const overlay = document.getElementById('settingsOverlay');
  const list = document.getElementById('settingsList');
  document.getElementById('settingsError').style.display = 'none';
  document.getElementById('settingsSuccess').style.display = 'none';
  overlay.classList.add('active');
  list.innerHTML = '<div style="text-align:center;color:var(--muted);padding:20px;">Loading categories…</div>';
  try {
    const res = await fetch('/settings');
    const data = await res.json();
    currentSettings = data.settings;
    const st = document.getElementById('smartToggle');
    if (st) st.checked = (currentSettings.redaction_style === 'smart');
    let html = '';
    for (const [groupName, groupData] of Object.entries(currentSettings.categories)) {
      const meta = GROUP_META[groupName] || { label: groupName, cls: '' };
      const col = GROUP_COLOR[groupName] || '#64748b';
      const subs = Object.entries(groupData.subcategories);
      let rows = '';
      subs.forEach(([catKey, catData]) => {
        const checked = catData.enabled ? 'checked' : '';
        const crit = catData.critical ? '<span class="crit">CRITICAL</span>' : '';
        rows += '<div class="cat-row" data-cat="' + catKey + '" data-label="' + (catData.description||catKey).toLowerCase() + '">'
          + '<label for="cb_' + catKey + '">' + escapeHtml(catData.description) + crit + '</label>'
          + '<span class="switch"><input type="checkbox" id="cb_' + catKey + '" ' + checked
          + ' onchange="toggleCategory(\'' + groupName + '\',\'' + catKey + '\',this.checked)"><span class="slider"></span></span></div>';
      });
      html += '<div class="group ' + meta.cls + '" data-group="' + groupName + '">'
        + '<div class="group-head" onclick="this.parentElement.classList.toggle(\'collapsed\')">'
        + '<div class="group-title"><span class="group-dot" style="background:' + col + '"></span>' + meta.label + ' <span class="group-count">(' + subs.length + ')</span></div>'
        + '<div class="group-tools"><button class="mini-btn" onclick="event.stopPropagation();groupAll(\'' + groupName + '\',true)">All</button>'
        + '<button class="mini-btn" onclick="event.stopPropagation();groupAll(\'' + groupName + '\',false)">None</button>'
        + '<span class="chev">▼</span></div></div>'
        + '<div class="group-body">' + rows + '</div></div>';
    }
    if (currentSettings.custom && currentSettings.custom.length) {
      let c = '';
      currentSettings.custom.forEach((custom, idx) => {
        const checked = custom.enabled ? 'checked' : '';
        c += '<div class="custom-item"><span>' + escapeHtml(custom.name) + ' — ' + escapeHtml(custom.description||'') + '</span>'
          + '<span class="switch"><input type="checkbox" ' + checked + ' onchange="toggleCustom(' + idx + ',this.checked)"><span class="slider"></span></span></div>';
      });
      html += '<div class="group"><div class="group-head"><div class="group-title"><span class="group-dot" style="background:#0ea5e9"></span>Custom Categories</div></div><div class="group-body">' + c + '</div></div>';
    }
    list.innerHTML = html;
    filterCats();
  } catch(e) {
    document.getElementById('settingsError').textContent = 'Error loading settings: ' + e.message;
    document.getElementById('settingsError').style.display = 'block';
  }
}
function filterCats(){
  const q = (document.getElementById('catSearch').value || '').toLowerCase().trim();
  document.querySelectorAll('#settingsList .group').forEach(g => {
    let any = false;
    g.querySelectorAll('.cat-row').forEach(r => {
      const show = !q || r.getAttribute('data-label').includes(q) || r.getAttribute('data-cat').toLowerCase().includes(q);
      r.style.display = show ? '' : 'none';
      if (show) any = true;
    });
    g.style.display = any ? '' : 'none';
  });
}
function groupAll(group, enabled){
  if (!currentSettings || !currentSettings.categories[group]) return;
  Object.keys(currentSettings.categories[group].subcategories).forEach(k => {
    currentSettings.categories[group].subcategories[k].enabled = enabled;
  });
  openSettings();
}
function setAll(enabled){
  if (!currentSettings) return;
  Object.values(currentSettings.categories).forEach(g => Object.values(g.subcategories).forEach(s => s.enabled = enabled));
  (currentSettings.custom||[]).forEach(c => c.enabled = enabled);
  openSettings();
}
function setStyle(on){
  if (currentSettings) currentSettings.redaction_style = on ? 'smart' : 'token';
}
function toggleCategory(group, catKey, enabled){
  if (currentSettings && currentSettings.categories[group] && currentSettings.categories[group].subcategories[catKey])
    currentSettings.categories[group].subcategories[catKey].enabled = enabled;
}
function toggleCustom(index, enabled){
  if (currentSettings && currentSettings.custom && currentSettings.custom[index])
    currentSettings.custom[index].enabled = enabled;
}
async function addCustomCategory(){
  const name = document.getElementById('customName').value.trim();
  const pattern = document.getElementById('customPattern').value.trim();
  const desc = document.getElementById('customDesc').value.trim() || name;
  const prefix = document.getElementById('customPrefix').value;
  const err = document.getElementById('settingsError');
  if (!name || !pattern) { err.textContent = 'Please enter both category name and regex pattern'; err.style.display = 'block'; return; }
  if (!currentSettings.custom) currentSettings.custom = [];
  currentSettings.custom.push({ name, pattern, description: desc, dummy_prefix: prefix, enabled: true });
  document.getElementById('customName').value = '';
  document.getElementById('customPattern').value = '';
  document.getElementById('customDesc').value = '';
  openSettings();
}
async function saveSettings(){
  try {
    const res = await fetch('/settings', { method: 'POST', headers: { 'Content-Type': 'application/json' }, body: JSON.stringify(currentSettings) });
    const data = await res.json();
    if (res.ok) {
      document.getElementById('settingsSuccess').style.display = 'block';
      setTimeout(() => { document.getElementById('settingsSuccess').style.display = 'none'; closeSettings(); }, 1400);
    } else { const e = document.getElementById('settingsError'); e.textContent = data.error || 'Save failed'; e.style.display = 'block'; }
  } catch(e) { const er = document.getElementById('settingsError'); er.textContent = 'Error: ' + e.message; er.style.display = 'block'; }
}
function closeSettings(){ document.getElementById('settingsOverlay').classList.remove('active'); }

loadDocs();
setInterval(loadDocs, 10000);
</script>
</body>
</html>
"""

# ─── API Server ───────────────────────────────────────────────────────────────
def main():
    port = 8765
    if "--api-server" in sys.argv:
        idx = sys.argv.index("--api-server")
        port = int(sys.argv[idx + 1]) if idx + 1 < len(sys.argv) else 8765

    settings = load_settings()
    engine = EnhancedRedactionEngine(settings=settings)

    def process_file(filepath, original_filename=None):
        """Process a single file - returns safe doc and redaction map."""
        doc_id = f"doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        original_name = original_filename or Path(filepath).name
        original_text = extract_text(filepath)

        # Redact PII
        style = engine.settings.get("redaction_style", "token") if engine.settings else "token"
        result = engine.redact(original_text, style=style)
        safe_text, redaction_map, category_counts, redactions, real_to_dummy = result

        # Convert list-based redaction_map to enriched format
        # redaction_map is now a list of dicts: [{variable, original, category, group, position}, ...]
        enriched_map = OrderedDict()
        for entry in redaction_map:
            var_name = entry["variable"]
            category = entry["category"]
            group_info = None
            for group_name, categories in SECURITY_POLICY.items():
                if category in categories:
                    group_info = group_name
                    break
            
            enriched_map[var_name] = {
                "original": entry["original"],
                "category": category,
                "group": group_info or "UNKNOWN",
                "description": entry.get("group", ""),
                "position": entry.get("position", 0)
            }

        # Save redaction map locally (never exposed via API)
        # In smart mode, the meaningful reversible mapping is real_value -> dummy_value.
        # In token mode, it is token -> original. We store BOTH for full fidelity.
        if real_to_dummy:
            simple_map = dict(real_to_dummy)                       # real -> dummy
            reverse_map = {v: k for k, v in real_to_dummy.items()}  # dummy -> real
        else:
            simple_map = {k: v["original"] for k, v in enriched_map.items()}  # token -> original
            reverse_map = {v["original"]: k for k, v in enriched_map.items()}  # original -> token
        map_file = MAPS_DIR / f"{doc_id}_redaction_map.json"
        map_data = {
            "document_id": doc_id,
            "original_filename": original_name,
            "redaction_style": style,
            "map": simple_map,           # smart: real->dummy | token: token->original
            "reverse_map": reverse_map,  # smart: dummy->real | token: original->token
            "full_map": enriched_map,    # Detailed mapping with categories
            "category_counts": category_counts,
            "created_at": datetime.now().isoformat(),
            "total_redactions": len(redactions)
        }
        map_file.write_text(json.dumps(map_data, indent=2))

        # Create enriched safe document
        # Note: redaction_map in safe doc contains ONLY variable->original mapping
        # The safe_text field contains the fully redacted text
        safe_doc = {
            "document_id": doc_id,
            "original_filename": original_name,
            "original_type": Path(filepath).suffix.lower(),
            "file_size": Path(filepath).stat().st_size,
            "processed_at": datetime.now().isoformat(),
            "total_redactions": len(redactions),
            "category_counts": category_counts,
            # SAFE OUTPUT: All sensitive data is now in variable form
            "all_text": safe_text,
            "processed_text": safe_text,
            "pages": [{"page_number": 1, "text": safe_text}],
            # Redaction map stored separately, NOT in safe output
            # to ensure zero PII leakage in the safe document
            "redaction_summary": {
                "total_categories": len(category_counts),
                "categories": list(category_counts.keys()),
                "total_variables": len(enriched_map)
            },
            "metadata": {"encoding": "utf-8"}
        }

        # Save safe doc
        safe_file = DOCS_DIR / f"{doc_id}_safe.json"
        safe_file.write_text(json.dumps(safe_doc, indent=2, ensure_ascii=False))

        return safe_doc, map_data, str(filepath), str(safe_file), str(map_file)

    def save_doc(safe_doc, rmap):
        doc_id = safe_doc["document_id"]
        safe_file = DOCS_DIR / f"{doc_id}_safe.json"
        safe_file.write_text(json.dumps(safe_doc, indent=2, ensure_ascii=False))

    class Handler(BaseHTTPRequestHandler):
        def log_message(self, fmt, *args): pass

        def do_GET(self):
            parsed = urlparse(self.path)
            path = parsed.path

            if path == "/health":
                self._json(200, {"status": "ok", "service": "OAKAI Document Reader"})
            elif path in ("/", "/ui", "/index.html"):
                self.send_response(200)
                self.send_header("Content-Type", "text/html; charset=utf-8")
                self.end_headers()
                self.write_html(HTML_UI)
            elif path == "/settings":
                # Return current settings
                current_settings = load_settings()
                # Also return policy info for new categories
                self._json(200, {
                    "settings": current_settings,
                    "policy": {
                        group: {
                            cat_key: {
                                "description": cat_info.get("description", cat_key),
                                "dummy_prefix": cat_info.get("dummy_prefix", cat_key)
                            }
                            for cat_key, cat_info in categories.items()
                        }
                        for group, categories in SECURITY_POLICY.items()
                    }
                })
            elif path == "/documents":
                docs = []
                if DOCS_DIR.exists():
                    for f in sorted(DOCS_DIR.glob("*_safe.json"), reverse=True)[:10]:
                        doc_id = f.stem.replace("_safe", "")
                        doc_data = json.loads(f.read_text())
                        docs.append({
                            "id": doc_id,
                            "filename": doc_data.get("original_filename", f.name),
                            "url": f"/documents/{doc_id}/safe",
                            "size": f.stat().st_size,
                            "redactions_count": doc_data.get("total_redactions", 0)
                        })
                self._json(200, {"documents": docs, "count": len(docs)})
            else:
                match = re.match(r'/documents/([^/]+)/safe', path)
                if match:
                    doc_id = match.group(1)
                    safe_file = DOCS_DIR / f"{doc_id}_safe.json"
                    if safe_file.exists():
                        data = json.loads(safe_file.read_text())
                        self._json(200, data)
                    else:
                        self._json(404, {"error": f"Document {doc_id} not found"})
                else:
                    # Also support /documents/<id>/map endpoint
                    map_match = re.match(r'/documents/([^/]+)/map', path)
                    if map_match:
                        doc_id = map_match.group(1)
                        map_file = MAPS_DIR / f"{doc_id}_redaction_map.json"
                        if map_file.exists():
                            data = json.loads(map_file.read_text())
                            self._json(200, data)
                        else:
                            self._json(404, {"error": f"Map for document {doc_id} not found"})
                    else:
                        self._json(404, {"error": "Not found"})

        def do_POST(self):
            parsed = urlparse(self.path)
            path = parsed.path

            if path == "/settings":
                # Update settings from POST body
                try:
                    content_length = int(self.headers.get('Content-Length', 0))
                    body = self.rfile.read(content_length)
                    incoming = json.loads(body)
                    # Always merge over a FRESH default settings object so the full
                    # category tree is preserved even if the client sends a partial
                    # payload (e.g. only {"redaction_style": ...}). Never clobber.
                    defaults = get_default_settings()
                    merged = defaults
                    if isinstance(incoming, dict):
                        # Top-level scalars (redaction_style) first
                        if "redaction_style" in incoming and incoming["redaction_style"] in ("token", "smart"):
                            merged["redaction_style"] = incoming["redaction_style"]
                        # Categories: overlay enabled flags / critical / custom from incoming
                        inc_cats = incoming.get("categories", {}) if isinstance(incoming.get("categories"), dict) else {}
                        if inc_cats:
                            for g, gd in inc_cats.items():
                                if g not in merged["categories"]:
                                    merged["categories"][g] = {"enabled": True, "subcategories": {}}
                                if isinstance(gd, dict):
                                    if "enabled" in gd:
                                        merged["categories"][g]["enabled"] = bool(gd["enabled"])
                                    inc_subs = gd.get("subcategories", {}) if isinstance(gd.get("subcategories"), dict) else {}
                                    for ck, cv in inc_subs.items():
                                        if ck in merged["categories"][g]["subcategories"] and isinstance(cv, dict):
                                            if "enabled" in cv:
                                                merged["categories"][g]["subcategories"][ck]["enabled"] = bool(cv["enabled"])
                                            if "critical" in cv:
                                                merged["categories"][g]["subcategories"][ck]["critical"] = bool(cv["critical"])
                        # Custom categories
                        if isinstance(incoming.get("custom"), list):
                            merged["custom"] = incoming["custom"]
                    save_settings(merged)
                    # Rebuild engine with merged settings (always has full category tree)
                    engine.settings = merged
                    globals()["settings"] = merged
                    engine._build_patterns()
                    self._json(200, {"status": "ok", "message": "Settings saved"})
                except Exception as e:
                    self._json(500, {"error": str(e)})
            elif path == "/settings/categories":
                # Return all available categories from policy for settings UI
                categories_list = []
                for group_name, categories in SECURITY_POLICY.items():
                    for cat_key, cat_info in categories.items():
                        categories_list.append({
                            "group": group_name,
                            "key": cat_key,
                            "label": cat_info.get("description", cat_key),
                            "dummy_prefix": cat_info.get("dummy_prefix", cat_key),
                            "critical": cat_info.get("critical", False)
                        })
                self._json(200, {"categories": categories_list})
            elif path == "/upload":
                self._handle_upload()
            elif path == "/process":
                self._handle_process()
            elif path == "/restart":
                self._json(200, {
                    "status": "restarting",
                    "message": "Server is restarting. This page will reload automatically in a few seconds."
                })

                def _do_restart():
                    time.sleep(0.4)  # let the HTTP response finish flushing to the client
                    try:
                        if os.name == "nt":
                            # Windows: let the VBS helper kill + relaunch silently
                            helper = SCRIPT_DIR / "restart_helper.vbs"
                            subprocess.Popen(
                                ["wscript.exe", str(helper)],
                                creationflags=subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP,
                                close_fds=True,
                            )
                        else:
                            # Linux / macOS / WSL2: re-exec this script in place so
                            # the very same process image (and thus the new code) runs
                            os.execv(sys.executable, [sys.executable, str(SCRIPT_DIR / "doc_reader_onefile.py")])
                    except Exception:
                        pass
                    os._exit(0)

                threading.Thread(target=_do_restart, daemon=True).start()
            else:
                self._json(404, {"error": "Not found"})

        def _handle_upload(self):
            """Handle multipart file upload."""
            try:
                content_type = self.headers.get("Content-Type", "")
                if not content_type.startswith("multipart/form-data"):
                    self._json(400, {"error": "Expected multipart/form-data"})
                    return

                boundary = content_type.split("boundary=")[1].encode()
                body = self.rfile.read(int(self.headers["Content-Length"]))
                parts = body.split(b"--" + boundary)

                for part in parts:
                    if b"filename=" in part:
                        # Extract filename
                        header_end = part.find(b"\r\n\r\n")
                        headers_raw = part[:header_end].decode("utf-8", errors="replace")
                        filename_match = re.search(r'filename="(.*?)"', headers_raw)
                        filename = filename_match.group(1) if filename_match else "unknown"

                        # Extract file data
                        file_data = part[header_end + 4:]
                        if file_data.endswith(b"\r\n"):
                            file_data = file_data[:-2]

                        if filename:
                            server_name = f"upload_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}"
                            server_path = UPLOAD_DIR / server_name
                            server_path.write_bytes(file_data)

                            safe_doc, rmap, saved_path, safe_path, map_path = process_file(
                                str(server_path), filename
                            )
                            self._json(200, {
                                "document_id": safe_doc["document_id"],
                                "safe_url": f"/documents/{safe_doc['document_id']}/safe",
                                "original_filename": filename,
                                "total_redactions": safe_doc["total_redactions"],
                                "category_counts": safe_doc["category_counts"],
                                "file_path": saved_path,
                                "safe_path": safe_path,
                                "map_path": map_path
                            })
                            return

                self._json(400, {"error": "No file uploaded"})
            except Exception as e:
                self._json(500, {"error": str(e)})

        def _handle_process(self):
            """Handle process request with file path."""
            try:
                body = self.rfile.read(int(self.headers["Content-Length"]))
                data = json.loads(body)
                fp = data.get("file_path", "")

                if not fp or not Path(fp).exists():
                    self._json(400, {"error": f"File not found: {fp}"})
                    return

                safe_doc, rmap, saved_path, safe_path, map_path = process_file(fp)
                self._json(200, {
                    "document_id": safe_doc["document_id"],
                    "safe_url": f"/documents/{safe_doc['document_id']}/safe",
                    "total_redactions": safe_doc["total_redactions"],
                    "category_counts": safe_doc["category_counts"],
                    "file_path": saved_path,
                    "safe_path": safe_path,
                    "map_path": map_path
                })
            except Exception as e:
                self._json(500, {"error": str(e)})

        def write_html(self, content):
            self.wfile.write(content.encode("utf-8"))

        def _json(self, code, data):
            self.send_response(code)
            self.send_header("Content-Type", "application/json")
            self.end_headers()
            self.wfile.write(json.dumps(data).encode())

    # Brief retry loop: if this process was just launched by restart_helper.vbs, the
    # previous instance may not have fully released the port yet. Retrying for a few
    # seconds avoids a silent bind failure racing against that cleanup.
    server = None
    last_err = None
    for attempt in range(10):
        try:
            server = HTTPServer(("0.0.0.0", port), Handler)
            break
        except OSError as e:
            last_err = e
            time.sleep(0.5)
    if server is None:
        print(f"❌ Could not bind port {port} after retries: {last_err}")
        sys.exit(1)
    print(f"✅ OAKAI Document Reader Server Running!")
    print(f"   Access via: http://localhost:{port}")
    print(f"   Health:     http://localhost:{port}/health")
    print(f"   UI:         http://localhost:{port}/")
    print(f"   Uploads:    {UPLOAD_DIR}")
    print(f"   Safe docs:  {DOCS_DIR}")
    print(f"   Maps:       {MAPS_DIR}")
    server.serve_forever()

if __name__ == "__main__":
    main()